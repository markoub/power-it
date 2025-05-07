"""
Analyze a PowerPoint file to understand its structure
"""
from pptx import Presentation
import sys
import os

def analyze_pptx(pptx_path):
    print(f"Analyzing PowerPoint file: {pptx_path}")
    print(f"File exists: {os.path.exists(pptx_path)}")
    
    p = Presentation(pptx_path)
    
    print(f"\nTotal slides: {len(p.slides)}")
    print(f"Total layouts: {len(p.slide_layouts)}")
    
    print("\nLayouts:")
    for i, layout in enumerate(p.slide_layouts):
        name = getattr(layout, 'name', f'Layout {i}')
        print(f"  Layout {i}: '{name}'")
        
        # Print placeholder info
        for j, placeholder in enumerate(layout.placeholders):
            ph_type = placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else 'unknown'
            ph_idx = placeholder.placeholder_format.idx if hasattr(placeholder, 'placeholder_format') else 'unknown'
            ph_name = getattr(placeholder, 'name', f'Placeholder {j}')
            print(f"    Placeholder {j}: '{ph_name}', Type: {ph_type}, Index: {ph_idx}")
    
    print("\nSlides:")
    for i, slide in enumerate(p.slides):
        print(f"\nSlide {i+1}:")
        
        # Get slide layout name if available
        layout_idx = -1
        for j, layout in enumerate(p.slide_layouts):
            if slide.slide_layout == layout:
                layout_idx = j
                break
        
        layout_name = p.slide_layouts[layout_idx].name if layout_idx >= 0 else "Unknown"
        print(f"  Layout: {layout_name} (index {layout_idx})")
        
        # Print shapes info
        print(f"  Shapes: {len(slide.shapes)}")
        for j, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, 'name', f'Shape {j}')
            
            # Try to get placeholder info if it's a placeholder
            placeholder_info = ""
            if hasattr(shape, 'placeholder_format'):
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                placeholder_info = f", Placeholder Type: {ph_type}, Index: {ph_idx}"
            
            # Get shape text if it's a text frame
            shape_text = ""
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                text = shape.text_frame.text
                # Truncate long text
                if len(text) > 30:
                    text = text[:27] + '...'
                if text:
                    shape_text = f", Text: '{text}'"
            
            print(f"    Shape {j}: '{shape_name}'{placeholder_info}{shape_text}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        pptx_path = "../test_toc.pptx"
    
    analyze_pptx(pptx_path) 