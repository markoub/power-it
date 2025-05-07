from pptx import Presentation
from tools.pptx_shapes import find_shape_by_name
from tools.pptx_toc import _find_section_text_shapes, process_toc_slide

# Load the template file
prs = Presentation("template.pptx")

# Print all the slide layouts
print("All slide layouts:")
for idx, layout in enumerate(prs.slide_layouts):
    layout_name = getattr(layout, 'name', f'Layout {idx}')
    print(f"  {idx}: {layout_name}")

# Try to find the TableOfContents layout
toc_layout = None
for idx, layout in enumerate(prs.slide_layouts):
    layout_name = getattr(layout, 'name', f'Layout {idx}')
    if 'tableofcontents' in layout_name.lower().replace(' ', ''):
        print(f"\nFound TableOfContents layout at index {idx}")
        toc_layout = layout
        break

if not toc_layout:
    print("\nTableOfContents layout not found specifically by name!")
    # Fallback to index 12 if available (as done in the code)
    if len(prs.slide_layouts) > 12:
        print(f"Falling back to layout 12: {getattr(prs.slide_layouts[12], 'name', 'Layout 12')}")
        toc_layout = prs.slide_layouts[12]
    else:
        print("Not enough layouts, can't fallback to layout 12")

if toc_layout:
    # Create a slide using the TOC layout
    test_slide = prs.slides.add_slide(toc_layout)
    
    print("\nAll shapes in the test slide created from TOC layout:")
    for i, shape in enumerate(test_slide.shapes):
        shape_name = getattr(shape, "name", f"Shape {i} (no name)")
        shape_text = ""
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            shape_text = shape.text_frame.text
        print(f"  {i}: {shape_name} - Text: {shape_text}")
    
    # Try our updated _find_section_text_shapes function
    section_shapes = _find_section_text_shapes(test_slide)
    
    print("\nSection shapes found by _find_section_text_shapes on test slide:")
    if not section_shapes:
        print("  No section shapes found!")
    else:
        for section_num, shape in section_shapes.items():
            shape_text = ""
            if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                shape_text = shape.text_frame.text
            print(f"  Section {section_num}: {shape.name} - Text: {shape_text}")
    
    # Test updating the sections
    test_sections = [f"Updated Section {i}" for i in range(1, 9)]
    process_toc_slide(test_slide, test_sections)
    
    print("\nAfter updating section texts:")
    for i, shape in enumerate(test_slide.shapes):
        shape_name = getattr(shape, "name", f"Shape {i} (no name)")
        shape_text = ""
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            shape_text = shape.text_frame.text
        print(f"  {i}: {shape_name} - Text: {shape_text}")
    
    # Don't save the temporary presentation (would overwrite template)
    # Just exit 