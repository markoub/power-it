"""
Template analyzer for PowerPoint presentations
Examines the structure of a template PPTX file and outputs details about layouts and shapes
"""
import os
import sys
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_layout(prs, layout_name=None, layout_index=None):
    """Analyze a specific layout by name or index"""
    target_layout = None
    
    # First, try to find the layout by name
    if layout_name:
        for i, layout in enumerate(prs.slide_layouts):
            current_name = getattr(layout, 'name', f'Layout {i}')
            if layout_name.lower() in current_name.lower():
                target_layout = layout
                print(f"\n===== DETAILED ANALYSIS OF '{current_name}' LAYOUT (index {i}) =====")
                break
    
    # If not found by name and index provided, use index
    if target_layout is None and layout_index is not None:
        if 0 <= layout_index < len(prs.slide_layouts):
            target_layout = prs.slide_layouts[layout_index]
            layout_name = getattr(target_layout, 'name', f'Layout {layout_index}')
            print(f"\n===== DETAILED ANALYSIS OF '{layout_name}' LAYOUT (index {layout_index}) =====")
    
    if target_layout is None:
        print(f"Layout not found: '{layout_name}' or index {layout_index}")
        return
    
    # Create a slide with this layout for analysis
    slide = prs.slides.add_slide(target_layout)
    
    # General layout information
    print("\nLAYOUT GENERAL INFO:")
    print(f"- Name: {getattr(target_layout, 'name', 'Unknown')}")
    print(f"- Total placeholders: {len(target_layout.placeholders)}")
    print(f"- Total shapes: {len(slide.shapes)}")
    
    # Analyze placeholders
    print("\nPLACEHOLDERS IN THIS LAYOUT:")
    placeholders = []
    for i, placeholder in enumerate(target_layout.placeholders):
        ph_type = placeholder.placeholder_format.type
        ph_name = placeholder.name
        ph_idx = placeholder.placeholder_format.idx
        
        # Try to get text if available
        ph_text = ""
        if hasattr(placeholder, "text") and placeholder.text:
            ph_text = f"Text: '{placeholder.text[:50]}...'" if len(placeholder.text) > 50 else f"Text: '{placeholder.text}'"
        
        # Get dimensions if available
        ph_dimensions = ""
        if hasattr(placeholder, "width") and hasattr(placeholder, "height"):
            ph_dimensions = f"Dimensions: ({placeholder.width}, {placeholder.height})"
        
        placeholders.append({
            "index": i,
            "name": ph_name,
            "ph_type": ph_type,
            "ph_idx": ph_idx,
            "text": ph_text,
            "dimensions": ph_dimensions
        })
        
        print(f"  Placeholder {i}: Name='{ph_name}', Type={ph_type}, Index={ph_idx}")
        if ph_text:
            print(f"    {ph_text}")
        if ph_dimensions:
            print(f"    {ph_dimensions}")
    
    # Analyze all shapes in the created slide
    print("\nSHAPES IN THE CREATED SLIDE:")
    for i, shape in enumerate(slide.shapes):
        shape_name = getattr(shape, "name", "Unknown")
        shape_type = type(shape).__name__
        shape_mso_type = "Unknown"
        
        try:
            # Try to get the MSO shape type
            shape_mso_type = shape.shape_type if hasattr(shape, "shape_type") else "Unknown"
        except:
            pass
        
        # Get shape text if available
        shape_text = ""
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            shape_text = shape.text_frame.text[:50] + "..." if len(shape.text_frame.text) > 50 else shape.text_frame.text
            
        # Check for placeholder properties
        placeholder_info = ""
        if hasattr(shape, "placeholder_format"):
            try:
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                placeholder_info = f"Placeholder Type={ph_type}, Index={ph_idx}"
            except:
                placeholder_info = "Has placeholder_format but couldn't access properties"
        
        # Get position and size
        pos_info = ""
        size_info = ""
        if hasattr(shape, "left") and hasattr(shape, "top"):
            pos_info = f"Position: ({shape.left}, {shape.top})"
        if hasattr(shape, "width") and hasattr(shape, "height"):
            size_info = f"Size: ({shape.width}, {shape.height})"
        
        print(f"\n  Shape {i}: Name='{shape_name}'")
        print(f"    Type: {shape_type}, MSO_Type: {shape_mso_type}")
        if shape_text:
            print(f"    Text: '{shape_text}'")
        if placeholder_info:
            print(f"    {placeholder_info}")
        if pos_info:
            print(f"    {pos_info}")
        if size_info:
            print(f"    {size_info}")
        
        # For text frames, analyze paragraphs and runs
        if hasattr(shape, "text_frame"):
            try:
                print(f"    Text Frame Analysis:")
                print(f"      Auto-size: {shape.text_frame.auto_size if hasattr(shape.text_frame, 'auto_size') else 'Unknown'}")
                print(f"      Word-wrap: {shape.text_frame.word_wrap if hasattr(shape.text_frame, 'word_wrap') else 'Unknown'}")
                print(f"      Paragraphs: {len(shape.text_frame.paragraphs)}")
                
                for j, para in enumerate(shape.text_frame.paragraphs):
                    print(f"      Paragraph {j}:")
                    print(f"        Text: '{para.text}'")
                    print(f"        Level: {para.level}")
                    print(f"        Alignment: {para.alignment if hasattr(para, 'alignment') else 'Unknown'}")
                    print(f"        Runs: {len(para.runs)}")
                    
                    for k, run in enumerate(para.runs):
                        print(f"          Run {k}: '{run.text}'")
                        if hasattr(run, "font"):
                            print(f"            Font: Size={run.font.size if hasattr(run.font, 'size') else 'Default'}, "
                                  f"Bold={run.font.bold if hasattr(run.font, 'bold') else 'Default'}, "
                                  f"Italic={run.font.italic if hasattr(run.font, 'italic') else 'Default'}")
            except Exception as e:
                print(f"    Error analyzing text frame: {str(e)}")
    
    print("\n=== LAYOUT ANALYSIS COMPLETE ===\n")
    return target_layout, slide

def analyze_template(template_path, layout_name=None, layout_index=None):
    """Analyze a PowerPoint template file and print details about layouts and shapes"""
    print(f"\n===== Analyzing template: {template_path} =====")
    print(f"Template exists: {os.path.exists(template_path)}")
    
    prs = Presentation(template_path)
    
    # If a specific layout is requested, analyze only that
    if layout_name or layout_index is not None:
        return analyze_layout(prs, layout_name, layout_index)
    
    # Otherwise, analyze all layouts
    print(f"\nFound {len(prs.slide_layouts)} slide layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        layout_name = getattr(layout, 'name', f'Layout {i}')
        placeholder_count = len(layout.placeholders)
        print(f"\nLayout {i}: '{layout_name}' ({placeholder_count} placeholders)")
        
        # Analyze placeholders in this layout
        for j, placeholder in enumerate(layout.placeholders):
            ph_type = placeholder.placeholder_format.type
            ph_name = placeholder.name
            ph_idx = placeholder.placeholder_format.idx
            
            # Try to get text if available
            ph_text = ""
            if hasattr(placeholder, "text") and placeholder.text:
                ph_text = f" - Text: '{placeholder.text[:30]}...'" if len(placeholder.text) > 30 else f" - Text: '{placeholder.text}'"
            
            print(f"  Placeholder {j}: Name='{ph_name}', Type={ph_type}, Index={ph_idx}{ph_text}")
    
    return None, None

def compare_layouts(prs, layout1_name, layout2_name):
    """Compare two layouts to identify differences"""
    print(f"\n===== COMPARING LAYOUTS: '{layout1_name}' and '{layout2_name}' =====")
    
    layout1, slide1 = analyze_layout(prs, layout1_name)
    layout2, slide2 = analyze_layout(prs, layout2_name)
    
    if not layout1 or not layout2:
        print("Cannot compare layouts: one or both layouts not found")
        return
    
    print("\n===== COMPARISON RESULTS =====")
    print(f"Layout 1: '{getattr(layout1, 'name', 'Unknown')}' - {len(layout1.placeholders)} placeholders, {len(slide1.shapes)} shapes")
    print(f"Layout 2: '{getattr(layout2, 'name', 'Unknown')}' - {len(layout2.placeholders)} placeholders, {len(slide2.shapes)} shapes")
    
    # Compare placeholder counts and types
    print("\nPlaceholder Type Comparison:")
    placeholders1 = [(p.placeholder_format.type, p.placeholder_format.idx, p.name) for p in layout1.placeholders]
    placeholders2 = [(p.placeholder_format.type, p.placeholder_format.idx, p.name) for p in layout2.placeholders]
    
    types1 = [p[0] for p in placeholders1]
    types2 = [p[0] for p in placeholders2]
    
    common_types = set(types1).intersection(set(types2))
    unique_to_1 = set(types1) - set(types2)
    unique_to_2 = set(types2) - set(types1)
    
    print(f"Common placeholder types: {common_types}")
    print(f"Placeholder types unique to '{getattr(layout1, 'name', 'Unknown')}': {unique_to_1}")
    print(f"Placeholder types unique to '{getattr(layout2, 'name', 'Unknown')}': {unique_to_2}")
    
    print("\nPlaceholder Index Comparison:")
    indices1 = [p[1] for p in placeholders1]
    indices2 = [p[1] for p in placeholders2]
    
    common_indices = set(indices1).intersection(set(indices2))
    print(f"Common placeholder indices: {common_indices}")
    print(f"Placeholder indices unique to '{getattr(layout1, 'name', 'Unknown')}': {set(indices1) - set(indices2)}")
    print(f"Placeholder indices unique to '{getattr(layout2, 'name', 'Unknown')}': {set(indices2) - set(indices1)}")
    
    # Compare shapes
    print("\nShape Type Comparison:")
    shape_types1 = [type(s).__name__ for s in slide1.shapes]
    shape_types2 = [type(s).__name__ for s in slide2.shapes]
    
    print(f"'{getattr(layout1, 'name', 'Unknown')}' shape types: {set(shape_types1)}")
    print(f"'{getattr(layout2, 'name', 'Unknown')}' shape types: {set(shape_types2)}")
    
    print("\n===== COMPARISON COMPLETE =====\n")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Analyze PowerPoint template file')
    parser.add_argument('template_file', nargs='?', default='template.pptx', help='Path to the template file')
    parser.add_argument('--layout-name', help='Name of layout to analyze in detail')
    parser.add_argument('--layout-index', type=int, help='Index of layout to analyze in detail')
    parser.add_argument('--compare', nargs=2, metavar=('LAYOUT1', 'LAYOUT2'), help='Compare two layouts')
    
    args = parser.parse_args()
    
    prs = Presentation(args.template_file)
    
    if args.compare:
        compare_layouts(prs, args.compare[0], args.compare[1])
    else:
        analyze_template(args.template_file, args.layout_name, args.layout_index) 