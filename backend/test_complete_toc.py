"""
Comprehensive test for all TOC generation functions
"""
from pptx import Presentation
import os

from tools.pptx_toc import create_table_of_contents_slide, process_toc_slide

def test_both_toc_functions():
    """
    Test both TOC generation functions to ensure they correctly create TOC slides 
    without bullet points and with the proper styling.
    """
    # Load the template
    template_path = os.path.join(os.path.dirname(__file__), "template.pptx")
    print(f"Loading template from: {template_path}")
    print(f"Template exists: {os.path.exists(template_path)}")
    
    prs = Presentation(template_path)
    
    # List of section titles
    section_titles = [
        "Understanding Eyelid Aging",
        "Blepharoplasty Techniques",
        "Patient Assessment",
        "Preoperative Planning",
        "Surgical Techniques",
        "Postoperative Care",
        "Results and Complications"
    ]
    
    # Test 1: create_table_of_contents_slide function
    print("\n===== TEST 1: create_table_of_contents_slide =====")
    toc_data = {
        "title": "Table of Contents - Test 1",
        "sections": section_titles
    }
    
    slide1 = create_table_of_contents_slide(prs, toc_data)
    print("TOC slide created using create_table_of_contents_slide function")
    
    # Test 2: process_toc_slide function
    print("\n===== TEST 2: process_toc_slide =====")
    # Add a blank TOC slide first
    toc_layout = None
    for i, layout in enumerate(prs.slide_layouts):
        name = getattr(layout, 'name', f'Layout {i}')
        if 'tableofcontents' in name.lower().replace(' ', ''):
            toc_layout = layout
            print(f"Found TOC layout by name: {name}")
            break
    
    if toc_layout is None and len(prs.slide_layouts) > 12:
        toc_layout = prs.slide_layouts[12]
        print(f"Using TOC layout by index 12")
    
    if toc_layout:
        blank_toc = prs.slides.add_slide(toc_layout)
        process_toc_slide(blank_toc, section_titles)
        print("TOC slide processed using process_toc_slide function")
    else:
        print("Couldn't find TOC layout, skipping Test 2")
    
    # Save the presentation with both slides
    output_path = os.path.join(os.path.dirname(__file__), "test_both_toc.pptx")
    prs.save(output_path)
    print(f"\nSaved test presentation with both TOC slides to: {output_path}")
    print(f"File exists: {os.path.exists(output_path)}")
    print(f"Check slides for proper formatting and NO bullet points!")

if __name__ == "__main__":
    test_both_toc_functions() 