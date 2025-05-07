"""
Test TOC slide generation
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import os

def test_toc_generation():
    # Load the template
    template_path = os.path.join(os.path.dirname(__file__), "template.pptx")
    print(f"Loading template from: {template_path}")
    print(f"Template exists: {os.path.exists(template_path)}")
    
    prs = Presentation(template_path)
    
    # Find the Table of Contents layout
    toc_layout = None
    for i, layout in enumerate(prs.slide_layouts):
        name = getattr(layout, 'name', f'Layout {i}')
        if 'tableofcontents' in name.lower().replace(' ', ''):
            toc_layout = layout
            print(f"Found TOC layout by name: {name}")
            break
    
    # If not found by name, use layout by index
    if toc_layout is None and len(prs.slide_layouts) > 12:
        toc_layout = prs.slide_layouts[12]
        print(f"Using TOC layout by index 12")
    
    # Fallback to content layout if still not found
    if toc_layout is None:
        for i, layout in enumerate(prs.slide_layouts):
            name = getattr(layout, 'name', f'Layout {i}')
            if 'content' in name.lower():
                toc_layout = layout
                print(f"Using content layout as fallback: {name}")
                break
    
    # Final fallback
    if toc_layout is None:
        toc_layout = prs.slide_layouts[1]  # Title and Content
        print("Using Title and Content layout as final fallback")
    
    # Create the TOC slide
    toc_slide = prs.slides.add_slide(toc_layout)
    
    # List of section titles
    section_titles = [
        "Understanding Eyelid Aging",
        "Blepharoplasty Techniques",
        "Patient Assessment",
        "Preoperative Planning",
        "Surgical Techniques",
        "Postoperative Care",
        "Results and Complications"
    ]
    
    # Print TOC slide shapes
    print("\nShapes in TOC slide:")
    for i, shape in enumerate(toc_slide.shapes):
        shape_name = getattr(shape, 'name', f'Shape {i}')
        placeholder_info = ""
        if hasattr(shape, 'placeholder_format'):
            placeholder_info = f", Placeholder Index: {shape.placeholder_format.idx}"
        print(f"  Shape {i}: '{shape_name}'{placeholder_info}")
    
    # Find title shape
    title_shape = None
    for shape in toc_slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # TITLE
            title_shape = shape
            break
    
    # Set title
    if title_shape and hasattr(title_shape, 'text_frame'):
        title_shape.text_frame.text = "Table of Contents"
        print(f"Set title to: Table of Contents")
    
    # Section indices mapping
    section_indices = {
        13: 1,  # First section uses placeholder with index 13
        14: 2,  # Second section uses placeholder with index 14 
        15: 3,
        16: 4,
        17: 5,
        18: 6,
        19: 7,
        20: 8
    }
    
    # Find section placeholders and populate them
    placed_sections = set()
    
    # First approach: Use placeholders based on indices
    print("\nPlacing sections by placeholder index:")
    for shape in toc_slide.shapes:
        if hasattr(shape, 'placeholder_format'):
            try:
                ph_idx = shape.placeholder_format.idx
                section_num = section_indices.get(ph_idx)
                
                if section_num and section_num <= len(section_titles):
                    section_title = section_titles[section_num-1]
                    
                    if hasattr(shape, 'text_frame'):
                        # Clear text frame
                        if hasattr(shape.text_frame, 'clear'):
                            shape.text_frame.clear()
                        else:
                            shape.text_frame.text = ""
                        
                        # Add content
                        if len(shape.text_frame.paragraphs) > 0:
                            paragraph = shape.text_frame.paragraphs[0]
                            paragraph.text = section_title
                            
                            # CRITICAL: Remove bullet points
                            if hasattr(paragraph, 'level'):
                                paragraph.level = 0
                                print(f"  Set section {section_num} paragraph level to 0 (no bullets)")
                            
                            # Keep alignment
                            if hasattr(paragraph, 'alignment'):
                                paragraph.alignment = PP_ALIGN.LEFT
                        else:
                            shape.text_frame.text = section_title
                            paragraph = shape.text_frame.paragraphs[0]
                            if hasattr(paragraph, 'level'):
                                paragraph.level = 0
                        
                        print(f"  Set section {section_num} placeholder (index {ph_idx}) to: '{section_title}'")
                        placed_sections.add(section_num)
            except Exception as e:
                print(f"  Error with placeholder {ph_idx}: {str(e)}")
    
    # Report results
    print(f"\nPlaced {len(placed_sections)} of {len(section_titles)} sections")
    
    # Second approach: Try content placeholder if needed
    if len(placed_sections) < len(section_titles):
        print("\nNot all sections placed, trying content placeholder")
        
        content_shape = None
        for shape in toc_slide.shapes:
            if (hasattr(shape, 'placeholder_format') and 
                shape.placeholder_format.type == 7 and  # OBJECT
                shape is not title_shape):
                content_shape = shape
                break
        
        if content_shape and hasattr(content_shape, 'text_frame'):
            content_shape.text_frame.clear()
            
            for section_num in range(1, len(section_titles) + 1):
                if section_num in placed_sections:
                    continue  # Skip sections already placed
                    
                section_title = section_titles[section_num-1]
                
                if section_num == min(section_num for section_num in range(1, len(section_titles) + 1) 
                                     if section_num not in placed_sections):
                    # First unplaced section
                    paragraph = content_shape.text_frame.paragraphs[0]
                else:
                    paragraph = content_shape.text_frame.add_paragraph()
                
                paragraph.text = section_title
                
                # CRITICAL: Remove bullet points
                if hasattr(paragraph, 'level'):
                    paragraph.level = 0
                    print(f"  Set content section {section_num} paragraph level to 0 (no bullets)")
                
                # Basic formatting
                if hasattr(paragraph, 'alignment'):
                    paragraph.alignment = PP_ALIGN.LEFT
                
                print(f"  Added section {section_num} to content placeholder: '{section_title}'")
    
    # Save the presentation
    output_path = os.path.join(os.path.dirname(__file__), "test_output_toc.pptx")
    prs.save(output_path)
    print(f"\nSaved test TOC presentation to: {output_path}")
    print(f"File exists: {os.path.exists(output_path)}")

if __name__ == "__main__":
    test_toc_generation() 