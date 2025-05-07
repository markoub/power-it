"""
Test file for validating Table of Contents (TOC) slide generation
"""

import os
import asyncio
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from models import SlidePresentation, Slide

from tools.generate_pptx import generate_pptx_from_slides

async def test_toc_generation():
    """Test the generation of a Table of Contents slide"""
    print("\n===== Testing TOC Generation =====")
    
    # Create test data from the provided JSON
    test_slides_data = [
        {
            "type": "Welcome",
            "fields": {
                "title": "Eyelid Rejuvenation: Surgical & Non-Surgical Approaches",
                "subtitle": "An Overview for Surgeons & Anti-Aging Professionals",
                "author": "Ana Malivukovic"
            }
        },
        {
            "type": "TableOfContents",
            "fields": {
                "title": "Table of Contents",
                "sections": [
                    "Understanding Eyelid Aging",
                    "Surgical Approaches",
                    "Non-Surgical Options",
                    "Treatment Algorithms"
                ]
            }
        },
        {
            "type": "Section",
            "fields": {
                "title": "Understanding Eyelid Aging"
            }
        },
        {
            "type": "Content",
            "fields": {
                "title": "The Aging Eyelid: More Than Just Cosmetics",
                "content": [
                    "Eyes as Focal Point: The eyelids often show aging first, causing a tired or aged appearance. Common complaints include droopy upper lids, under-eye bags, dark circles, and fine wrinkles.",
                    "Aesthetic vs. Functional Issues: Beyond cosmetics, severe upper lid drooping (dermatochalasis) can impair vision, and lower lid laxity can cause dryness or irritation.",
                    "Prevalence: Eyelid lifts (blepharoplasty) are among the top 5 cosmetic surgeries worldwide, reflecting high demand. Patients increasingly seek less invasive treatments to avoid surgical risks.",
                    "Goal: Achieve a youthful, rested eye appearance safely and effectively, addressing skin, muscle, and fat changes with an individualized plan."
                ]
            }
        }
    ]
    
    # Convert to SlidePresentation object
    test_slides = []
    for slide_data in test_slides_data:
        test_slides.append(Slide(
            type=slide_data['type'],
            fields=slide_data['fields']
        ))
    
    slide_presentation = SlidePresentation(
        title="Eyelid Rejuvenation",
        slides=test_slides
    )
    
    # Generate PowerPoint presentation
    result = await generate_pptx_from_slides(slide_presentation, presentation_id="test")
    
    # Validate that the file was created
    assert os.path.exists(result.pptx_path), f"PPTX file not created at {result.pptx_path}"
    print(f"PPTX file created: {result.pptx_path}")
    
    # Analyze the generated PPTX to check if TOC was created properly
    prs = Presentation(result.pptx_path)
    
    print(f"\nTotal slides in presentation: {len(prs.slides)}")
    
    # Look for TOC slide (should be the second slide)
    toc_slide = prs.slides[1] if len(prs.slides) > 1 else None
    
    if toc_slide:
        print(f"\nTOC slide found with {len(toc_slide.shapes)} shapes")
        
        # Print details about each shape in the TOC slide
        for i, shape in enumerate(toc_slide.shapes):
            shape_name = getattr(shape, "name", "Unknown")
            shape_type = type(shape).__name__
            
            # Check if shape has text
            shape_text = ""
            if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                shape_text = f" - Text: '{shape.text_frame.text[:30]}...'" if len(shape.text_frame.text) > 30 else f" - Text: '{shape.text_frame.text}'"
            
            # Check for placeholder info
            placeholder_info = ""
            if hasattr(shape, "placeholder_format"):
                try:
                    ph_type = shape.placeholder_format.type
                    ph_idx = shape.placeholder_format.idx
                    placeholder_info = f" - Placeholder Type={ph_type}, Index={ph_idx}"
                except:
                    placeholder_info = " - Has placeholder_format but couldn't access properties"
            
            print(f"  Shape {i}: Name='{shape_name}', Type={shape_type}{shape_text}{placeholder_info}")
            
            # Check if any of the shapes contain the section titles
            if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                for section in test_slides_data[1]['fields']['sections']:
                    if section in shape.text_frame.text:
                        print(f"    FOUND SECTION: '{section}' in shape {i}")
    else:
        print("\nTOC slide not found!")
    
    print("\n===== TOC Test Complete =====")
    return result.pptx_path

if __name__ == "__main__":
    asyncio.run(test_toc_generation()) 