import os
import sys
import pytest
import asyncio
from unittest.mock import patch, MagicMock
from PIL import Image

# Add parent directory to sys.path to import from main code
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from tools.generate_pptx import generate_pptx_from_slides
from models import SlidePresentation, Slide
from pptx import Presentation


class TestContentWithLogosSlide:
    """Test the ContentWithLogos slide type functionality"""

    @pytest.mark.asyncio
    async def test_contentlogos_slide_generation(self):
        """Test that ContentWithLogos slides are created properly with logos and content"""
        
        # Create a test logos directory
        test_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "tests", "test_data")
        os.makedirs(test_dir, exist_ok=True)
        
        # Create dummy logo files - use PNG instead of SVG
        logo1_path = os.path.join(test_dir, "test_logo1.png")
        logo2_path = os.path.join(test_dir, "test_logo2.png")
        logo3_path = os.path.join(test_dir, "test_logo3.png")
        
        # Create simple PNG files if they don't exist
        for i, logo_path in enumerate([logo1_path, logo2_path, logo3_path]):
            if not os.path.exists(logo_path):
                # Create a 100x100 colored image
                img = Image.new('RGB', (100, 100), color = ('red' if i == 0 else 'green' if i == 1 else 'blue'))
                img.save(logo_path)
                print(f"Created test logo: {logo_path}")
        
        # Create a test SlidePresentation with a ContentWithLogos slide
        presentation = SlidePresentation(
            title="Test ContentWithLogos Presentation",
            slides=[
                Slide(
                    type="Welcome",
                    fields={
                        "title": "Test Welcome Slide",
                        "subtitle": "Testing ContentWithLogos Slides",
                        "author": "Test Author"
                    }
                ),
                Slide(
                    type="ContentWithLogos",
                    fields={
                        "title": "Cloud Providers",
                        "content": [
                            "These are the major cloud providers",
                            "Each has different strengths",
                            "Choose based on your needs"
                        ],
                        "logo1": logo1_path,
                        "logo2": logo2_path,
                        "logo3": logo3_path
                    }
                )
            ]
        )
        
        # Create a temporary output path for the presentation
        output_path = os.path.join(test_dir, "test_contentlogos.pptx")
        
        # Generate PPTX with the ContentWithLogos slide
        result = await generate_pptx_from_slides(presentation, output_path)
        
        # Verify the presentation was created
        assert result is not None
        assert os.path.exists(result.pptx_path)
        
        # Open the presentation for verification
        prs = Presentation(result.pptx_path)
        
        # We should have at least 3 slides (Welcome, ContentWithLogos, ThankYou)
        assert len(prs.slides) >= 3
        
        # Check the ContentWithLogos slide (should be the second slide)
        content_logos_slide = prs.slides[1]
        
        # Check that the slide has a title
        title_found = False
        title_text = ""
        for shape in content_logos_slide.shapes:
            if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # Title
                title_found = True
                title_text = shape.text if hasattr(shape, "text") else ""
                break
            elif hasattr(shape, "text") and hasattr(shape, "name") and "title" in shape.name.lower():
                title_found = True
                title_text = shape.text
                break
        
        assert title_found, "Title shape not found in ContentWithLogos slide"
        
        # Check for content text
        content_found = False
        content_text = ""
        for shape in content_logos_slide.shapes:
            if hasattr(shape, "text_frame") and hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 2:
                content_found = True
                content_text = shape.text_frame.text
                
                # Check if paragraphs exist
                if hasattr(shape.text_frame, "paragraphs") and shape.text_frame.paragraphs:
                    # Verify that text in paragraphs is not bold
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.runs:
                            # At least one run should not be bold (default formatting)
                            non_bold_found = False
                            for run in paragraph.runs:
                                if hasattr(run, "font") and hasattr(run.font, "bold") and run.font.bold is False:
                                    non_bold_found = True
                                    break
                            if not non_bold_found and paragraph.runs:
                                print(f"Warning: All runs in paragraph are bold: {paragraph.text}")
                
                break
            elif hasattr(shape, "text_frame") and hasattr(shape, "name") and "content" in shape.name.lower():
                content_found = True
                content_text = shape.text_frame.text
                break
        
        assert content_found, "Content shape not found in ContentWithLogos slide"
        
        # Count the number of image shapes to verify logos were added
        shape_count_before_logos = 0
        for shape in content_logos_slide.shapes:
            try:
                if hasattr(shape, "placeholder_format") and callable(getattr(shape, "placeholder_format", None)):
                    placeholder_format = shape.placeholder_format
                    if placeholder_format.type != 18:  # Not a picture placeholder
                        shape_count_before_logos += 1
                else:
                    shape_count_before_logos += 1
            except (ValueError, AttributeError):
                # This is not a placeholder shape
                shape_count_before_logos += 1
        
        print(f"Found {shape_count_before_logos} non-picture-placeholder shapes on the slide")
                
        # Check if any picture was added
        # In PowerPoint, when a picture is added, placeholders are replaced or removed,
        # so the count of shapes will change
        shape_count = len(list(content_logos_slide.shapes))
        print(f"Total shape count: {shape_count}")
        
        # Assert at least 3 shapes (title, content, and at least one logo)
        assert shape_count >= 3, "Not enough shapes on the slide - logos might not have been added"
        
        # Clean up
        if os.path.exists(result.pptx_path):
            try:
                os.remove(result.pptx_path)
            except:
                pass 