import os
import pytest
import tempfile
import shutil
from pptx import Presentation

# Add the parent directory to the path so we can import the module
import sys
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from tools.pptx_slides import create_section_slide
from tools.pptx_utils import format_section_number

def test_real_section_slide():
    """Test that the section slide is created correctly with proper title and number using the real template."""
    # Get the template file path
    template_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "template.pptx")
    assert os.path.exists(template_path), f"Template file does not exist at {template_path}"
    
    # Load the template
    prs = Presentation(template_path)
    
    # Find the section layout
    section_layout = None
    for layout in prs.slide_layouts:
        if hasattr(layout, 'name') and layout.name == 'Section':
            section_layout = layout
            break
    
    assert section_layout is not None, "Section layout not found in template"
    
    # Create a temporary directory for the output PPTX
    temp_dir = tempfile.mkdtemp()
    try:
        # Create a section slide
        section_title = "Test Section"
        section_number = 2
        slide_obj = create_section_slide(prs, section_layout, section_title, section_number)
        
        # Save the presentation
        output_pptx = os.path.join(temp_dir, "test_section_slide.pptx")
        prs.save(output_pptx)
        
        # Verify the file was created
        assert os.path.exists(output_pptx), f"Output file not created at {output_pptx}"
        
        # Load the saved file and check the section slide
        prs2 = Presentation(output_pptx)
        slide = prs2.slides[0]  # The section slide should be the first slide
        
        # Verify shapes have the correct text
        title_found = False
        number_found = False
        expected_number = format_section_number(section_number)
        
        # Print all shapes for debugging
        print("\nShapes in the section slide:")
        for i, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, "name", f"Shape {i}")
            shape_text = ""
            placeholder_type = "N/A"
            
            if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                shape_text = shape.text_frame.text
                
            if hasattr(shape, "placeholder_format"):
                placeholder_type = shape.placeholder_format.type
                
            print(f"  Shape {i}: Name='{shape_name}', Type={placeholder_type}, Text='{shape_text}'")
            
            # Check if this is the title or number shape
            if shape_text == section_title:
                title_found = True
                print(f"  Found section title: '{shape_text}'")
                
            if shape_text == expected_number:
                number_found = True
                print(f"  Found section number: '{shape_text}'")
        
        # Final verification
        assert title_found, f"Section title '{section_title}' not found in any shape"
        assert number_found, f"Section number '{expected_number}' not found in any shape"
        
    finally:
        # Clean up
        shutil.rmtree(temp_dir)

if __name__ == "__main__":
    test_real_section_slide() 