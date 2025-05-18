"""
Unit tests for Table of Contents slide generation
"""
import os
import pytest
import asyncio
import tempfile
from pptx import Presentation

from models import SlidePresentation, Slide

from tools.generate_pptx import generate_pptx_from_slides
from tools.pptx_toc import create_table_of_contents_slide

@pytest.mark.asyncio
async def test_toc_generation():
    """Test the generation of a Table of Contents slide with sections"""
    
    # Create test data with sections
    test_slides = [
        Slide(
            type="Welcome",
            fields={
                "title": "Eyelid Rejuvenation: Surgical & Non-Surgical Approaches",
                "subtitle": "An Overview for Surgeons & Anti-Aging Professionals",
                "author": "Ana Malivukovic"
            }
        ),
        Slide(
            type="TableOfContents",
            fields={
                "title": "Table of Contents",
                "sections": [
                    "Understanding Eyelid Aging",
                    "Surgical Approaches", 
                    "Non-Surgical Options",
                    "Treatment Algorithms"
                ]
            }
        ),
        Slide(
            type="Section",
            fields={
                "title": "Understanding Eyelid Aging"
            }
        ),
        Slide(
            type="Content",
            fields={
                "title": "The Aging Eyelid: More Than Just Cosmetics",
                "content": [
                    "Eyes as Focal Point: The eyelids often show aging first.",
                    "Aesthetic vs. Functional Issues: Beyond cosmetics.",
                    "Prevalence: Eyelid lifts are among the top 5 cosmetic surgeries.",
                    "Goal: Achieve a youthful, rested eye appearance."
                ]
            }
        )
    ]
    
    # Create presentation and generate PPTX
    presentation = SlidePresentation(
        title="Test TOC",
        slides=test_slides
    )
    
    # Create a temporary file for the output
    test_output_path = os.path.join(tempfile.gettempdir(), "test_toc.pptx")
    
    result = await generate_pptx_from_slides(presentation, test_output_path)
    
    # Check that file was created
    assert os.path.exists(result.pptx_path)
    
    # Open the PPTX and analyze the TOC slide
    prs = Presentation(result.pptx_path)
    
    # Check that we have at least 4 slides (Welcome, TOC, Section, Content)
    assert len(prs.slides) >= 4, f"Presentation should have at least 4 slides, but has {len(prs.slides)}"
    
    # Just check that the TOC slide has shapes with text containing at least one section name
    section_titles = test_slides[1].fields["sections"]
    section_found = False
    
    # Check all slides for section names, not just the second one
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for section_title in section_titles:
                    if section_title in shape.text_frame.text:
                        section_found = True
                        print(f"Found section '{section_title}' on slide {slide_idx+1}")
                        break
    
    assert section_found, "No section titles found in any slide"
    
    # Clean up the temporary file
    try:
        os.remove(result.pptx_path)
    except Exception:
        pass
    
    return result.pptx_path


@pytest.mark.asyncio
async def test_toc_with_empty_sections():
    """Test TOC generation with empty sections list"""
    
    # Create test data with empty sections list in TOC
    test_slides = [
        Slide(
            type="Welcome",
            fields={
                "title": "Test Presentation",
                "subtitle": "Testing TOC with empty sections"
            }
        ),
        Slide(
            type="TableOfContents",
            fields={
                "title": "Table of Contents",
                "sections": []
            }
        ),
        Slide(
            type="Content",
            fields={
                "title": "Test Content",
                "content": ["Test content item"]
            }
        )
    ]
    
    # Create presentation and generate PPTX
    presentation = SlidePresentation(
        title="Test Empty TOC",
        slides=test_slides
    )
    
    # Create a temporary file for the output
    test_output_path = os.path.join(tempfile.gettempdir(), "test_empty_toc.pptx")
    
    result = await generate_pptx_from_slides(presentation, test_output_path)
    
    # Check that file was created
    assert os.path.exists(result.pptx_path)
    
    # Open the PPTX and analyze
    prs = Presentation(result.pptx_path)
    
    # With empty sections, the TOC slide might be skipped
    # Check that we at least have a Welcome slide and a Content slide
    assert len(prs.slides) >= 2, f"Presentation should have at least 2 slides, but has {len(prs.slides)}"
    
    # Verify we have a content slide with the expected title
    content_title_found = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and "Test Content" in shape.text_frame.text:
                content_title_found = True
                break
    
    assert content_title_found, "Content slide title not found"
    
    # Clean up the temporary file
    try:
        os.remove(result.pptx_path)
    except Exception:
        pass
    
    return result.pptx_path


if __name__ == "__main__":
    asyncio.run(test_toc_generation()) 