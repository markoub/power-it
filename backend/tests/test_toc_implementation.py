"""
Test script to verify the Table of Contents implementation works with a real presentation.
"""
import pytest
from pptx import Presentation
from tools.pptx_toc import process_toc_slide, create_table_of_contents_slide
import os

@pytest.fixture
def template_presentation():
    """Fixture to provide a template presentation for tests."""
    template_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "template.pptx")
    return Presentation(template_path)

def test_create_toc_slide(template_presentation):
    """Test creating a new TOC slide."""
    toc_data = {
        'title': 'Table of Contents',
        'sections': [
            'Executive Summary',
            'Market Analysis',
            'Competitive Landscape',
            'Product Strategy',
            'Financial Projections',
            'Implementation Plan',
            'Risk Assessment',
            'Conclusion'
        ]
    }

    new_toc_slide = create_table_of_contents_slide(template_presentation, toc_data)
    
    # Verify the slide was created
    assert new_toc_slide is not None
    
    # Verify the slide has expected content
    found_sections = []
    for shape in new_toc_slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text
            for section in toc_data['sections']:
                if section in text:
                    found_sections.append(section)
    
    # Check that all sections are found in the slide
    for section in toc_data['sections']:
        assert section in found_sections, f"Section '{section}' not found in the TOC slide"

def test_update_toc_slide(template_presentation):
    """Test updating an existing TOC slide."""
    # First create a TOC slide
    initial_sections = [
        'Executive Summary',
        'Market Analysis',
        'Competitive Landscape',
    ]
    
    toc_data = {
        'title': 'Table of Contents',
        'sections': initial_sections
    }
    
    toc_slide = create_table_of_contents_slide(template_presentation, toc_data)
    
    # Then update it with new sections
    updated_sections = [
        'Updated Executive Summary',
        'Updated Market Analysis',
        'Updated Competitive Landscape',
        'New Section',
    ]
    
    process_toc_slide(toc_slide, updated_sections)
    
    # Verify the updated content
    found_sections = []
    for shape in toc_slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text
            for section in updated_sections:
                if section in text:
                    found_sections.append(section)
    
    # Check that all updated sections are found in the slide
    for section in updated_sections:
        assert section in found_sections, f"Updated section '{section}' not found in the TOC slide" 