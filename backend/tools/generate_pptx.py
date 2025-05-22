"""
Generate PowerPoint presentations from slide content
"""
import os
import uuid
import asyncio
import re
from typing import Dict, Any, List, Optional, Union

from pptx import Presentation

from config import PRESENTATIONS_STORAGE_DIR
from models import SlidePresentation, CompiledPresentation, Slide, CompiledSlide, PptxGeneration
from tools.slide_config import SLIDE_TYPES, PRESENTATION_STRUCTURE

# Import our refactored modules
from .pptx_utils import list_presentation_images
from .pptx_shapes import get_layout_by_name, find_shape_by_name, get_toc_shapes
from .pptx_text import adjust_text_size
from .pptx_toc import create_table_of_contents_slide, process_toc_slide

# Import all slide creation functions
from .pptx_slides import (
    create_welcome_slide, 
    create_section_slide, 
    create_content_slide,
    create_3images_slide,
    create_thank_you_slide,
    create_content_with_logos_slide
)

def format_section_number(number):
    """Format a section number to ensure it's a two-digit string (e.g., 1 -> '01')."""
    return f"{number:02d}"

async def generate_pptx_from_slides(slides, output_path) -> PptxGeneration:
    """
    Generate a PowerPoint presentation from slide content using a template.
    
    Args:
        slides: SlidePresentation or CompiledPresentation object containing slide content,
               or a list of slides
        output_path: Path where to save the presentation
        
    Returns:
        A PptxGeneration object with the path to the PPTX file
    """
    try:
        # Handle different input types
        if isinstance(slides, list):
            # We've received a list of slide dictionaries
            slides_list = slides
            presentation_title = "Generated Presentation"
            # Try to get title from first Welcome slide if available
            for slide in slides_list:
                if slide.get('type') == 'Welcome':
                    presentation_title = slide.get('title', presentation_title)
                    break
        else:
            # We have a SlidePresentation or CompiledPresentation object
            slides_list = slides.slides
            presentation_title = getattr(slides, 'title', "Generated Presentation")
        
        # Create presentation directories if they don't exist
        if output_path is not None:
            if output_path == "test":
                presentation_dir = os.path.join(PRESENTATIONS_STORAGE_DIR, "test")
            else:
                presentation_dir = os.path.join(PRESENTATIONS_STORAGE_DIR, str(output_path))
            
            # Create presentation directory and images subdirectory
            os.makedirs(presentation_dir, exist_ok=True)
            images_dir = os.path.join(presentation_dir, "images")
            os.makedirs(images_dir, exist_ok=True)
            print(f"Created/ensured directories: {presentation_dir} and {images_dir}")
        
        # Print summary of slides to be processed
        print(f"\n===== Processing Presentation: {presentation_title} =====")
        print(f"Output path: {output_path}")
        print(f"Total slides to process: {len(slides_list)}")
        
        # Print slide types summary
        slide_types = {}
        image_slides = 0
        for slide in slides_list:
            slide_type = getattr(slide, 'type', None)
            if not slide_type and hasattr(slide, 'get'):
                slide_type = slide.get('type')
            
            slide_types[slide_type] = slide_types.get(slide_type, 0) + 1
            
            # Count slides that should contain images
            if slide_type in ['ContentImage', '3Images', 'ContentWithLogos']:
                image_slides += 1
        
        print("Slide types summary:")
        for slide_type, count in slide_types.items():
            print(f"  - {slide_type}: {count} slides")
        print(f"Total slides that should contain images: {image_slides}\n")
        
        # Debug images directory structure
        print(f"\n===== Debugging Image Structure for Presentation {output_path} =====")
        available_images = list_presentation_images(output_path, PRESENTATIONS_STORAGE_DIR)
        if available_images:
            print(f"Available images: {len(available_images)}")
            # Print details for first 10 images
            for i, img_path in enumerate(available_images[:10]):
                if os.path.exists(img_path):
                    size = os.path.getsize(img_path)
                    print(f"  - {os.path.basename(img_path)} (size: {size} bytes)")
                else:
                    print(f"  - {os.path.basename(img_path)} (file does not exist)")
            
            if len(available_images) > 10:
                print(f"  ... and {len(available_images) - 10} more images")
        else:
            print("No images found for this presentation!")
        print(f"========================================================\n")
        
        # Load template and extract layouts
        template_path = os.path.join(os.path.dirname(__file__), '../template.pptx')
        print(f"Loading template from: {template_path}")
        print(f"Template exists: {os.path.exists(template_path)}")
        
        # Create presentation object
        prs = Presentation(template_path)
        
        # Display layouts
        print("\n===== TEMPLATE STRUCTURE =====")
        print(f"Total slide layouts: {len(prs.slide_layouts)}")
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = getattr(layout, 'name', f'Layout {i}')
            print(f"Layout {i}: '{layout_name}'")
        print("=============================\n")
        
        # Generate presentation with the first slide layout
        welcome_layout = get_layout_by_name(prs, "Welcome")
        if welcome_layout:
            print(f"Found welcome layout by name: {getattr(welcome_layout, 'name', 'Welcome')}")
            # Print placeholder details
            placeholders = [shape for shape in welcome_layout.placeholders]
            print(f"  Placeholders in welcome layout: {len(placeholders)}")
            for i, placeholder in enumerate(placeholders):
                print(f"  - Placeholder {i}: Name='{placeholder.name}', Type={placeholder.placeholder_format.type}")
        else:
            welcome_layout = prs.slide_layouts[0]
            print(f"No welcome layout found, using first layout: {getattr(welcome_layout, 'name', 'Layout 0')}")
            
        # Find content layout
        content_image_layout = get_layout_by_name(prs, "ContentImage")
        if content_image_layout:
            print(f"Found ContentImage layout by name: {getattr(content_image_layout, 'name', 'ContentImage')}")
            # Print placeholder details
            placeholders = [shape for shape in content_image_layout.placeholders]
            print(f"  Placeholders in this layout: {len(placeholders)}")
            for i, placeholder in enumerate(placeholders):
                print(f"  - Placeholder {i}: Name='{placeholder.name}', Type={placeholder.placeholder_format.type}")
        else:
            content_image_layout = prs.slide_layouts[5]  # Typically the 6th layout in PowerPoint templates
            print(f"No ContentImage layout found, using layout 5: {getattr(content_image_layout, 'name', 'Layout 5')}")
        
        # Find TOC layout
        table_of_contents_layout = get_layout_by_name(prs, "TableOfContents")
        if table_of_contents_layout:
            print(f"Found TableOfContents layout by name: {getattr(table_of_contents_layout, 'name', 'TableOfContents')}")
        else:
            # Try to find a good fallback
            table_of_contents_layout = prs.slide_layouts[12]
            print(f"Using TableOfContents layout from index 12 directly")
        
        # Check if we have at least 13 layouts
        if table_of_contents_layout is None:
            print("WARNING: Template doesn't have enough layouts, using fallback for TOC layout")
            # Try to find a content layout instead
            for i, layout in enumerate(prs.slide_layouts):
                layout_name = getattr(layout, 'name', f'Layout {i}').lower()
                if 'content' in layout_name:
                    table_of_contents_layout = layout
                    print(f"Using {getattr(layout, 'name', f'Layout {i}')} as TOC layout")
                    break
            
            # Final fallback
            if table_of_contents_layout is None:
                table_of_contents_layout = prs.slide_layouts[0]
                print(f"Using first layout as TOC fallback")
        
        content_layout = get_layout_by_name(prs, "Content") or prs.slide_layouts[5]
        section_layout = get_layout_by_name(prs, "Section") or prs.slide_layouts[2]
        thankyou_layout = get_layout_by_name(prs, "ThankYou") or prs.slide_layouts[14] 
        
        # If ThankYou layout not found, try to find a Blank layout or use the last layout
        if not thankyou_layout:
            for i, layout in enumerate(prs.slide_layouts):
                layout_name = getattr(layout, 'name', '')
                if 'blank' in layout_name.lower() or 'thank' in layout_name.lower():
                    thankyou_layout = layout
                    print(f"Using '{layout_name}' as ThankYou slide")
                    break
            
            if not thankyou_layout:
                # Fall back to the first blank-looking layout or just use index 0
                thankyou_layout = prs.slide_layouts[0]
                print(f"No ThankYou layout found, using first layout as fallback")
        
        print(f"\nGenerating PPTX presentation with {len(slides_list)} slides")
        
        # Add welcome slide as the first slide
        # Find welcome slide data from the slides collection
        welcome_data = {}
        for slide in slides_list:
            if getattr(slide, 'type', None) == 'Welcome' and hasattr(slide, 'fields'):
                welcome_data = slide.fields
                print(f"Found Welcome slide data: {welcome_data}")
                break
        
        welcome_slide = create_welcome_slide(prs, welcome_layout, presentation_title, welcome_data)
        
        # Check if we need to add a table of contents slide
        sections = [slide for slide in slides_list if getattr(slide, 'type', None) == 'Section']
        toc_data = None
        
        # Look for a TableOfContents type slide in the data
        for slide in slides_list:
            if getattr(slide, 'type', None) == 'TableOfContents' and hasattr(slide, 'fields'):
                toc_data = slide.fields
                print(f"\nFound TableOfContents slide data: {toc_data}")
                break
        
        # If we don't have toc_data, create it from section slides
        if not toc_data:
            toc_data = {
                "title": "Table of Contents",
                "sections": []
            }
            # Extract section titles from section slides
            for i, slide in enumerate(sections):
                if hasattr(slide, 'fields') and 'title' in slide.fields:
                    toc_data["sections"].append(slide.fields['title'])
                else:
                    toc_data["sections"].append(f"Section {i+1}")
            
            print(f"Created TOC data from sections: {toc_data}")
        
        # Use our dedicated function to create the TOC slide correctly
        if toc_data and toc_data.get("sections"):
            # Create a TOC slide with properly formatted textboxes
            toc_slide = create_table_of_contents_slide(prs, toc_data)
            print("TOC slide created successfully with section titles")
        else:
            print("No TOC data available, skipping TOC slide")
        
        # Track section count for numbering
        section_count = 0
        
        # Add all content slides
        for i, slide in enumerate(slides_list):
            # Get slide title from fields
            slide_title = slide.fields.get("title", f"Slide {i+1}") if hasattr(slide, "fields") else getattr(slide, "title", f"Slide {i+1}")
            slide_type = getattr(slide, 'type', 'ContentImage')
            print(f"\n===== Processing slide {i+1}: {slide_title} =====")
            print(f"SLIDE DEBUG: Slide type detected: '{slide_type}'")
            print(f"SLIDE DEBUG: Slide fields: {getattr(slide, 'fields', {}).keys()}")
            
            # Skip welcome and TOC slides since we've already created them
            if slide_type in ['Welcome', 'TableOfContents']:
                print(f"Skipping {slide_type} slide - already created")
                continue
            
            # Log slide type for debugging
            print(f"SLIDE DEBUG: Processing slide type: '{slide_type}'")
            
            # Get slide fields
            fields = getattr(slide, "fields", {})
            
            # Create appropriate slide based on type
            if slide_type == 'Section':
                section_count += 1
                create_section_slide(prs, section_layout, slide_title, section_count)
            elif slide_type == 'Content':
                content = fields.get("content", [])
                create_content_slide(prs, content_layout, slide_title, content)
            elif slide_type == 'ContentWithLogos':
                print(f"SLIDE DEBUG: Found ContentWithLogos slide! Title: {slide_title}")
                print(f"SLIDE DEBUG: Fields: {fields.keys()}")
                content = fields.get("content", [])
                
                print(f"\nLOGO DEBUG: === Processing ContentWithLogos slide {i} ===")
                print(f"LOGO DEBUG: Original fields for ContentWithLogos slide:")
                for k, v in fields.items():
                    if k.startswith('logo'):
                        print(f"LOGO DEBUG:   - {k}: {v}")
                
                # Collect logo paths
                logo_paths = {}
                
                # Import logo fetcher directly here to ensure it's available
                from tools.logo_fetcher import download_logo
                
                for logo_key in ['logo1', 'logo2', 'logo3']:
                    if logo_key in fields and fields[logo_key]:
                        logo_value = fields[logo_key]
                        print(f"LOGO DEBUG: Processing {logo_key}: {logo_value}")
                        
                        # Check if logo_value is a file path that exists
                        if os.path.exists(logo_value):
                            logo_paths[logo_key] = logo_value
                            print(f"LOGO DEBUG: Found existing logo file for {logo_key}: {logo_value}")
                        else:
                            # Extract a logo term for fetching
                            logo_term = logo_value
                            
                            # If it's a file path that doesn't exist, extract the company name
                            if '/' in logo_value:
                                # Extract file name without extension
                                logo_term = os.path.basename(logo_value)
                                logo_term = os.path.splitext(logo_term)[0]
                                # Clean up any prefixes like slide_id_ or logo_
                                logo_term = re.sub(r'^(slide_id_|slide_|logo_)', '', logo_term)
                                # Remove logo number suffix if present
                                logo_term = re.sub(r'_(logo\d+)$', '', logo_term)
                            
                            print(f"LOGO DEBUG: Attempting to download logo for term: '{logo_term}'")
                            try:
                                # Try to download the logo
                                success, result = download_logo(logo_term)
                                if success:
                                    logo_paths[logo_key] = result
                                    print(f"LOGO DEBUG: Successfully downloaded logo for {logo_key}: {result}")
                                    
                                    # Verify the file exists and has content
                                    if os.path.exists(result):
                                        file_size = os.path.getsize(result)
                                        print(f"LOGO DEBUG: Logo file size: {file_size} bytes")
                                        if file_size == 0:
                                            print(f"LOGO DEBUG: WARNING - Logo file is empty!")
                                    else:
                                        print(f"LOGO DEBUG: WARNING - Downloaded logo file doesn't exist!")
                                else:
                                    print(f"LOGO DEBUG: Failed to download logo for {logo_key}: {result}")
                                    
                                    # Try with a simpler term (just use the company name directly)
                                    simple_term = logo_term.split('_')[0]
                                    if simple_term != logo_term:
                                        print(f"LOGO DEBUG: Trying simpler term: '{simple_term}'")
                                        success, result = download_logo(simple_term)
                                        if success:
                                            logo_paths[logo_key] = result
                                            print(f"LOGO DEBUG: Successfully downloaded logo with simpler term: {result}")
                            except Exception as e:
                                print(f"LOGO DEBUG: Error downloading logo for {logo_key}: {str(e)}")
                                import traceback
                                print(f"LOGO DEBUG: Error details:\n{traceback.format_exc()}")
                
                # Count how many valid logo paths we have
                valid_logos = sum(1 for path in logo_paths.values() if path and os.path.exists(path))
                print(f"LOGO DEBUG: Found {valid_logos} valid logos for ContentWithLogos slide")
                
                if valid_logos == 0:
                    print(f"LOGO DEBUG: NO VALID LOGOS FOUND - Creating slide will not have logos!")
                
                # Find appropriate layout for ContentWithLogos
                logo_content_layout = None
                
                # First try to find by exact name
                for i, layout in enumerate(prs.slide_layouts):
                    layout_name = getattr(layout, 'name', '')
                    if layout_name == 'ContentWithLogos':
                        logo_content_layout = layout
                        print(f"Found ContentWithLogos layout by exact name match")
                        break
                
                # If not found by exact name, try partial match
                if not logo_content_layout:
                    for i, layout in enumerate(prs.slide_layouts):
                        layout_name = getattr(layout, 'name', f'Layout {i}').lower().replace(' ', '')
                        if 'logo' in layout_name or 'contentlogo' in layout_name:
                            logo_content_layout = layout
                            print(f"Found ContentWithLogos layout by partial name: {getattr(layout, 'name', f'Layout {i}')}")
                            break
                
                # If still not found, look for layouts with picture placeholders for logos
                if not logo_content_layout:
                    print("Looking for layouts with picture placeholders that could be used for logos")
                    for i, layout in enumerate(prs.slide_layouts):
                        # Check if this layout has picture placeholders
                        logo_placeholders = 0
                        for shape in layout.placeholders:
                            if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 18:  # PICTURE
                                logo_placeholders += 1
                        
                        # If we found at least one picture placeholder and a content placeholder,
                        # this might be suitable for ContentWithLogos
                        if logo_placeholders > 0:
                            has_content = any(hasattr(shape, "placeholder_format") and 
                                            shape.placeholder_format.type == 2
                                            for shape in layout.placeholders)
                            if has_content:
                                logo_content_layout = layout
                                print(f"Found suitable layout with {logo_placeholders} picture placeholders: {getattr(layout, 'name', f'Layout {i}')}")
                                break
                
                # If no specific layout found, use content layout as fallback
                if not logo_content_layout:
                    print("No ContentWithLogos layout found, using content layout as fallback")
                    logo_content_layout = content_layout
                
                create_content_with_logos_slide(prs, logo_content_layout, slide_title, content, logo_paths)
            elif slide_type == '3Images':
                # Find appropriate layout for 3Images
                threeim_layout = None
                for i, layout in enumerate(prs.slide_layouts):
                    layout_name = getattr(layout, 'name', f'Layout {i}').lower().replace(' ', '')
                    if '3images' in layout_name or 'threeimages' in layout_name or 'images3' in layout_name:
                        threeim_layout = layout
                        print(f"Found 3Images layout by name: {getattr(layout, 'name', f'Layout {i}')}")
                        break
                
                # If no specific 3Images layout found, use content layout as fallback
                if not threeim_layout:
                    print("No 3Images layout found, using content layout as fallback")
                    threeim_layout = content_layout
                
                # Get image paths and subtitles from fields
                image_paths = {
                    'image1': None,
                    'image2': None,
                    'image3': None
                }
                
                # Handle both naming conventions for subtitles (subtitleimage1 or image1subtitle)
                subtitles = {}
                for i in range(1, 4):
                    # Check both possible field name formats and use the first one found
                    subtitle_value = fields.get(f'image{i}subtitle', '') or fields.get(f'subtitleimage{i}', '')
                    subtitles[f'image{i}subtitle'] = subtitle_value
                
                # Check for image fields in the slide data
                if output_path is not None:
                    images_dir = os.path.join(PRESENTATIONS_STORAGE_DIR, str(output_path), "images")
                    print(f"\nLooking for 3Images slide images in: {images_dir}")
                    if os.path.exists(images_dir):
                        print(f"Images directory exists with files: {os.listdir(images_dir)}")
                        
                        # Try to find images by slide ID if available
                        slide_id = getattr(slide, 'id', None)
                        slide_index = i + 1  # 1-based for filenames
                        
                        for img_num in range(1, 4):
                            img_field = f'image{img_num}'
                            print(f"Looking for {img_field}...")
                            
                            # First try to find image by specific image field name
                            if hasattr(slide, 'fields') and img_field in slide.fields:
                                # Handle direct path or pattern
                                image_attr = slide.fields[img_field]
                                if image_attr:
                                    print(f"  Found {img_field} in slide fields: {image_attr}")
                                    
                                    # If it's a direct path that exists, use it
                                    if os.path.exists(image_attr):
                                        file_size = os.path.getsize(image_attr)
                                        if file_size > 0:
                                            image_paths[img_field] = image_attr
                                            print(f"  Using direct {img_field} path: {image_attr} (size: {file_size} bytes)")
                                            continue
                                        else:
                                            print(f"  Warning: Direct image file is empty: {image_attr}")
                            
                            # Try various naming patterns for images
                            image_patterns = [
                                # Pattern 1: slide_id_{slide_id}_{img_field}
                                f"slide_id_{slide_id}_{img_field}" if slide_id else None,
                                # Pattern 2: slide_{slide_index}_{img_field}
                                f"slide_{slide_index}_{img_field}",
                                # Pattern 3: just {img_field}
                                f"{img_field}"
                            ]
                            
                            found_image = False
                            for pattern in image_patterns:
                                if not pattern:
                                    continue
                                    
                                print(f"  Trying pattern: {pattern}")
                                matching_images = [f for f in os.listdir(images_dir) 
                                                  if pattern.lower() in f.lower()]
                                
                                if matching_images:
                                    # Print all matching images with timestamps for debugging
                                    print(f"  Found {len(matching_images)} images matching pattern '{pattern}':")
                                    for img in matching_images:
                                        img_path = os.path.join(images_dir, img)
                                        mod_time = os.path.getmtime(img_path)
                                        size = os.path.getsize(img_path)
                                        print(f"    - {img} (Modified: {mod_time}, Size: {size} bytes)")
                                    
                                    # Sort by modification time (newest first)
                                    matching_images.sort(key=lambda f: os.path.getmtime(os.path.join(images_dir, f)), 
                                                        reverse=True)
                                    
                                    # Try files until we find one that's not empty
                                    for img_file in matching_images:
                                        image_path = os.path.join(images_dir, img_file)
                                        file_size = os.path.getsize(image_path)
                                        
                                        if file_size > 0:
                                            image_paths[img_field] = image_path
                                            print(f"  Found {img_field} with pattern '{pattern}': {image_path} (size: {file_size} bytes)")
                                            found_image = True
                                            break
                                        else:
                                            print(f"  Warning: Image file is empty, trying next: {image_path}")
                                    
                                    if found_image:
                                        break
                            
                            if not found_image:
                                print(f"  No valid image found for {img_field}")
                                
                                # Last resort - look for any image with this number in the filename
                                number_pattern = f"{img_num}"
                                print(f"  Last resort - looking for any image with '{number_pattern}' in filename")
                                matching_images = [f for f in os.listdir(images_dir) 
                                                 if number_pattern in f and (
                                                     f.lower().endswith('.png') or 
                                                     f.lower().endswith('.jpg') or 
                                                     f.lower().endswith('.jpeg'))]
                                
                                if matching_images:
                                    # Take the newest non-empty image
                                    matching_images.sort(key=lambda f: os.path.getmtime(os.path.join(images_dir, f)), 
                                                       reverse=True)
                                    
                                    for img_file in matching_images:
                                        image_path = os.path.join(images_dir, img_file)
                                        file_size = os.path.getsize(image_path)
                                        
                                        if file_size > 0:
                                            image_paths[img_field] = image_path
                                            print(f"  Found {img_field} with number pattern: {image_path} (size: {file_size} bytes)")
                                            break
                
                # Validate all image paths with PIL if possible
                for key, path in image_paths.items():
                    if path:
                        try:
                            from PIL import Image
                            img = Image.open(path)
                            print(f"  Validated {key} using PIL - Format: {img.format}, Size: {img.size}")
                            img.close()
                        except Exception as e:
                            print(f"  Warning - Could not validate image for {key} with PIL: {str(e)}")
                            # Continue anyway, let the PPTX library handle it
                
                print(f"Final image paths for 3Images slide:")
                for key, path in image_paths.items():
                    if path:
                        file_exists = os.path.exists(path)
                        file_size = os.path.getsize(path) if file_exists else 0
                        print(f"  {key}: {path} (exists: {file_exists}, size: {file_size} bytes)")
                    else:
                        print(f"  {key}: Not found")
                
                # Create the 3Images slide
                create_3images_slide(prs, threeim_layout, slide_title, image_paths, subtitles)
                
            elif slide_type == 'ContentImage':
                content = fields.get("content", [])
                image_attr = fields.get('image')
                
                # Try to find the image file if we have output_path
                image_path = None
                
                # First, check if the slide has an image_url attribute (from compiled presentation)
                image_url = getattr(slide, 'image_url', None)
                if image_url:
                    print(f"  Using image URL from slide data: {image_url}")
                    # If it's a relative URL, convert to absolute path
                    if not image_url.startswith('http') and output_path is not None:
                        # Extract filename from URL path
                        image_filename = os.path.basename(image_url)
                        images_dir = os.path.join(PRESENTATIONS_STORAGE_DIR, str(output_path), "images")
                        if os.path.exists(images_dir):
                            # Try to find the exact file
                            if image_filename in os.listdir(images_dir):
                                image_path = os.path.join(images_dir, image_filename)
                                print(f"  Found exact image at: {image_path}")
                                
                                # Verify file is valid
                                if os.path.getsize(image_path) == 0:
                                    print(f"  WARNING: Image file is empty: {image_path}")
                                    image_path = None
                
                # First try to find by slide ID if available
                if not image_path and hasattr(slide, 'id') and output_path is not None:
                    slide_id = getattr(slide, 'id')
                    images_dir = os.path.join(PRESENTATIONS_STORAGE_DIR, str(output_path), "images")
                    print(f"  Looking for image by slide ID {slide_id} in {images_dir}")
                    
                    if os.path.exists(images_dir):
                        # Look for image with exact slide ID
                        id_pattern = f"slide_id_{slide_id}_"
                        matching_id_images = [f for f in os.listdir(images_dir) if id_pattern in f.lower()]
                        
                        if matching_id_images:
                            # Print all matching images with timestamps for debugging
                            print(f"  Found {len(matching_id_images)} images matching slide ID {slide_id}:")
                            for img in matching_id_images:
                                img_path = os.path.join(images_dir, img)
                                mod_time = os.path.getmtime(img_path)
                                size = os.path.getsize(img_path)
                                print(f"    - {img} (Modified: {mod_time}, Size: {size} bytes)")
                            
                            # Sort matching images by modification time (newest first)
                            matching_id_images.sort(key=lambda f: os.path.getmtime(os.path.join(images_dir, f)), reverse=True)
                            image_path = os.path.join(images_dir, matching_id_images[0])
                            print(f"  Found image by slide ID: {image_path} (Last modified: {os.path.getmtime(image_path)})")
                            
                            # Verify file is valid
                            if os.path.getsize(image_path) == 0:
                                print(f"  WARNING: Image file is empty: {image_path}")
                                # Try next image if available
                                if len(matching_id_images) > 1:
                                    image_path = os.path.join(images_dir, matching_id_images[1])
                                    print(f"  Trying next image instead: {image_path}")
                                else:
                                    image_path = None
                
                # Fall back to pattern matching only if no image_path was found
                if not image_path and output_path is not None:
                    images_dir = os.path.join(PRESENTATIONS_STORAGE_DIR, str(output_path), "images")
                    print(f"  Looking for images in directory: {images_dir}")
                    
                    if os.path.exists(images_dir):
                        # Find image matching slide pattern - index is 0-based internally but 1-based for filenames
                        # Use slide_index+1 to match slide_1_*, slide_2_*, etc.
                        slide_index = i + 1
                        image_pattern = f"slide_{slide_index}_"
                        print(f"  Looking for images with pattern: {image_pattern}")
                        matching_images = [f for f in os.listdir(images_dir) if image_pattern in f.lower()]
                        
                        # Also look for just 'image' in filename
                        if not matching_images:
                            print(f"  No images with slide index pattern found, looking for files with 'image' in the name")
                            matching_images = [f for f in os.listdir(images_dir) if 'image' in f.lower()]
                        
                        if matching_images:
                            # Print all matching images with their timestamps for debugging
                            print(f"  All matching images for slide {slide_index}:")
                            for img in matching_images:
                                img_path = os.path.join(images_dir, img)
                                mod_time = os.path.getmtime(img_path)
                                size = os.path.getsize(img_path)
                                print(f"    - {img} (Modified: {mod_time}, Size: {size} bytes)")
                            
                            # Sort matching images by modification time (newest first)
                            matching_images.sort(key=lambda f: os.path.getmtime(os.path.join(images_dir, f)), reverse=True)
                            image_path = os.path.join(images_dir, matching_images[0])
                            print(f"  Selected image for slide {slide_index}: {image_path} (Last modified: {os.path.getmtime(image_path)})")
                            
                            # Verify file is valid
                            if os.path.getsize(image_path) == 0:
                                print(f"  WARNING: Image file is empty: {image_path}")
                                # Try next image if available
                                if len(matching_images) > 1:
                                    image_path = os.path.join(images_dir, matching_images[1])
                                    print(f"  Trying next image instead: {image_path}")
                                else:
                                    image_path = None
                
                # Last resort - check if there's a direct path to an image file
                if not image_path and image_attr and os.path.exists(image_attr):
                    print(f"  Using direct image path from slide data: {image_attr}")
                    image_path = image_attr
                
                # Validate the image file using PIL if possible
                if image_path:
                    try:
                        from PIL import Image
                        img = Image.open(image_path)
                        print(f"  Successfully validated image with PIL - Format: {img.format}, Size: {img.size}")
                        img.close()
                    except Exception as e:
                        print(f"  Warning - Could not validate image with PIL: {str(e)}")
                        # Continue anyway, let the PPTX library handle it
                
                create_content_slide(
                    prs, 
                    content_image_layout, 
                    slide_title, 
                    content, 
                    add_image=bool(image_path), 
                    image_path=image_path
                )
            
            # Add speaker notes if available - check both old and new format
            notes = getattr(slide, 'notes', None) or fields.get('notes')
            if notes:
                # Get the slide we just added
                current_slide = prs.slides[-1]
                if hasattr(current_slide, "notes_slide") and current_slide.notes_slide:
                    current_slide.notes_slide.notes_text_frame.text = notes
                    print(f"  Added speaker notes")
        
        # Add ThankYou slide at the end
        create_thank_you_slide(prs, thankyou_layout)
        
        # Create directory for this presentation if it doesn't exist
        pptx_filename = f"presentation_{uuid.uuid4()}.pptx"
        
        if output_path is not None:
            if output_path == "test":
                presentation_dir = os.path.join(PRESENTATIONS_STORAGE_DIR, "test")
            else:
                presentation_dir = os.path.join(PRESENTATIONS_STORAGE_DIR, str(output_path))
            
            # Create presentation directory
            os.makedirs(presentation_dir, exist_ok=True)
            
            # Also create images directory to ensure it exists
            images_dir = os.path.join(presentation_dir, "images")
            os.makedirs(images_dir, exist_ok=True)
            print(f"Ensured directories exist: {presentation_dir} and {images_dir}")
            
            # Log existing images if any
            if os.path.exists(images_dir) and os.listdir(images_dir):
                print(f"Images directory contains: {os.listdir(images_dir)}")
            
            pptx_path = os.path.join(presentation_dir, pptx_filename)
        else:
            os.makedirs(os.path.join(PRESENTATIONS_STORAGE_DIR, "temp"), exist_ok=True)
            pptx_path = os.path.join(PRESENTATIONS_STORAGE_DIR, "temp", pptx_filename)
        
        # Save the presentation
        print(f"Saving PPTX to: {pptx_path}")
        prs.save(pptx_path)
        print(f"PPTX saved successfully: {os.path.exists(pptx_path)}")
        
        return PptxGeneration(
            presentation_id=output_path,
            pptx_filename=pptx_filename,
            pptx_path=pptx_path,
            slide_count=len(prs.slides)
        )

    except Exception as e:
        print(f"Error generating PPTX: {str(e)}")
        import traceback
        print(traceback.format_exc())
        raise

async def convert_pptx_to_png(pptx_path: str, output_dir: Optional[str] = None) -> List[str]:
    """Convert a PPTX file to PNG images using LibreOffice to first create a PDF, then convert PDF to PNGs."""
    print("DEBUG: Using backend/tools/generate_pptx.py implementation")
    if output_dir is None:
        output_dir = os.path.dirname(pptx_path)
    
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Get basename without extension
    basename = os.path.splitext(os.path.basename(pptx_path))[0]
    
    print(f"Converting PPTX file: {pptx_path}")
    print(f"Output directory: {output_dir}")
    
    # Step 1: Convert PPTX to PDF using LibreOffice
    pdf_path = os.path.join(output_dir, f"{basename}.pdf")
    cmd_to_pdf = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path]
    
    pdf_success = False
    try:
        print(f"Running command to create PDF: {' '.join(cmd_to_pdf)}")
        process = await asyncio.create_subprocess_exec(
            *cmd_to_pdf,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await process.communicate()

        stdout_text = stdout.decode() if stdout else 'None'
        stderr_text = stderr.decode() if stderr else 'None'
        
        print(f"LibreOffice PDF stdout: {stdout_text}")
        print(f"LibreOffice PDF stderr: {stderr_text}")
        print(f"LibreOffice PDF exit code: {process.returncode}")

        # Check if PDF was created successfully
        if process.returncode == 0 and os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
            pdf_success = True
            print(f"Successfully created PDF at: {pdf_path}")
        else:
            print(f"Failed to create PDF from PPTX")
    except Exception as e:
        print(f"Error in PDF creation: {str(e)}")
        import traceback
        print(traceback.format_exc())
    
    # Step 2: If PDF creation succeeded, convert PDF pages to PNGs
    png_files = []
    if pdf_success:
        try:
            # Try to use pdf2image if available (which uses poppler internally)
            try:
                from pdf2image import convert_from_path
                
                print(f"Converting PDF to PNG images using pdf2image")
                images = convert_from_path(pdf_path, dpi=150)
                
                print(f"Generated {len(images)} page images from PDF")
                
                # Save each image as PNG
                for i, image in enumerate(images):
                    png_filename = f"{basename}-{i+1:03d}.png"
                    png_path = os.path.join(output_dir, png_filename)
                    image.save(png_path, "PNG")
                    png_files.append(png_path)
                    print(f"Saved PNG for page {i+1}: {png_path}")
                
            except ImportError:
                print("pdf2image not available, trying PyMuPDF (fitz)")
                # Try PyMuPDF (fitz) as an alternative
                try:
                    import fitz  # PyMuPDF
                    
                    print(f"Converting PDF to PNG images using PyMuPDF")
                    pdf_document = fitz.open(pdf_path)
                    
                    for page_num in range(len(pdf_document)):
                        page = pdf_document.load_page(page_num)
                        
                        # Render page to an image (with higher resolution)
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        
                        # Save the image
                        png_filename = f"{basename}-{page_num+1:03d}.png"
                        png_path = os.path.join(output_dir, png_filename)
                        pix.save(png_path)
                        png_files.append(png_path)
                        print(f"Saved PNG for page {page_num+1}: {png_path}")
                    
                    pdf_document.close()
                except ImportError:
                    print("PyMuPDF not available, falling back to direct LibreOffice PNG conversion")
                    # If neither pdf2image nor PyMuPDF is available, fall back to direct PPTX to PNG conversion
                    pdf_success = False
                except Exception as e:
                    print(f"Error using PyMuPDF: {str(e)}")
                    pdf_success = False
                    
        except Exception as e:
            print(f"Error converting PDF to PNG: {str(e)}")
            import traceback
            print(traceback.format_exc())
            pdf_success = False
    
    # Step 3: If PDF conversion failed or we couldn't convert PDF to PNGs, try direct PPTX to PNG conversion
    if not pdf_success or not png_files:
        print("PDF method failed, falling back to direct PPTX to PNG conversion")
        
        # Use LibreOffice to export each slide to PNG images directly
        cmd = ["soffice", "--headless", "--convert-to", "png", "--outdir", output_dir, pptx_path]
        
        try:
            print(f"Running command for direct PNG conversion: {' '.join(cmd)}")
            process = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await process.communicate()

            stdout_text = stdout.decode() if stdout else 'None'
            stderr_text = stderr.decode() if stderr else 'None'
            
            print(f"LibreOffice PNG stdout: {stdout_text}")
            print(f"LibreOffice PNG stderr: {stderr_text}")
            print(f"LibreOffice PNG exit code: {process.returncode}")

            # LibreOffice converts PPTX files to a series of PNG files with specific naming patterns
            # First look for pattern: basename-001.png, basename-002.png, etc.
            pattern1 = f"{basename}-"
            # Alternate pattern sometimes used: basename.png
            pattern2 = f"{basename}.png"
            
            # Collect all PNG files in the directory
            all_png_files = [
                os.path.join(output_dir, f)
                for f in os.listdir(output_dir)
                if f.endswith('.png')
            ]
            
            print(f"Found {len(all_png_files)} total PNG files in directory")
            
            # Check for pattern matches
            pattern1_matches = [f for f in all_png_files if os.path.basename(f).startswith(pattern1)]
            pattern2_matches = [f for f in all_png_files if os.path.basename(f) == pattern2]
            
            print(f"Pattern '{pattern1}' matches: {len(pattern1_matches)}")
            print(f"Pattern '{pattern2}' matches: {len(pattern2_matches)}")
            
            # Use pattern1 if we have matches, otherwise try pattern2
            if pattern1_matches:
                png_files = pattern1_matches
                print(f"Using {len(png_files)} files matching pattern '{pattern1}'")
            elif pattern2_matches:
                png_files = pattern2_matches
                print(f"Using single file matching pattern '{pattern2}'")
            else:
                # If no specific matches found, use all recent PNG files created in the last 60 seconds
                # This is a fallback for cases where naming conventions might differ
                import time
                current_time = time.time()
                recent_png_files = [
                    f for f in all_png_files 
                    if os.path.getmtime(f) > current_time - 60  # files modified in the last minute
                ]
                
                if recent_png_files:
                    png_files = recent_png_files
                    print(f"Using {len(png_files)} recently created PNG files")
                else:
                    # Last resort: try to match PPTX filename in PNG files
                    partial_matches = [
                        f for f in all_png_files 
                        if any(part in os.path.basename(f).lower() for part in basename.lower().split('_'))
                    ]
                    
                    if partial_matches and len(partial_matches) > 1:  # Need more than 1 to have multiple slides
                        png_files = partial_matches
                        print(f"Using {len(png_files)} files with partial name match")
                    else:
                        # If still no matches, just use all PNG files as last resort
                        if len(all_png_files) > 1:  # Need more than 1 to have multiple slides
                            png_files = all_png_files
                            print(f"Using all {len(png_files)} PNG files as last resort")
                        else:
                            png_files = []

            # Sort files in a natural order (so -1, -2, -10 appear in the correct sequence)
            if png_files:
                png_files.sort(key=lambda path: [
                    int(part) if part.isdigit() else part
                    for part in re.split(r'(\d+)', path)
                ])

                print(f"Selected {len(png_files)} PNG files")
                for i, png_file in enumerate(png_files[:10]):  # Show just first 10 for brevity
                    print(f"  {i+1}. {os.path.basename(png_file)}")
                
                if len(png_files) > 10:
                    print(f"  ... and {len(png_files) - 10} more files")
            
        except Exception as e:
            print(f"Error in direct PPTX to PNG conversion: {str(e)}")
            import traceback
            print(traceback.format_exc())
    
    # Step 4: If we still don't have enough PNGs, generate them directly from the PPTX using python-pptx
    if not png_files or len(png_files) < 2:
        print("No PNG files found after conversion attempts. Generating them directly using python-pptx and PIL.")
        
        try:
            from pptx import Presentation
            from PIL import Image, ImageDraw, ImageFont
            
            # Load the presentation
            prs = Presentation(pptx_path)
            slide_count = len(prs.slides)
            
            print(f"Loaded presentation with {slide_count} slides")
            
            # Create PNGs for each slide
            png_files = []
            
            for i, slide in enumerate(prs.slides):
                slide_png_filename = f"{basename}-{i+1:03d}.png"
                slide_png_path = os.path.join(output_dir, slide_png_filename)
                
                # Create a blank image with a white background
                img = Image.new('RGB', (1280, 720), color=(255, 255, 255))
                d = ImageDraw.Draw(img)
                
                # Try to use a system font
                try:
                    font = ImageFont.truetype("Arial", 20)
                except IOError:
                    try:
                        font = ImageFont.truetype("DejaVuSans.ttf", 20)
                    except IOError:
                        font = ImageFont.load_default()
                
                # Add slide index and some text
                d.text((40, 40), f"Slide {i+1} of {slide_count}", fill=(0, 0, 0), font=font)
                
                # Try to extract title if present
                slide_title = None
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text:
                                slide_title = paragraph.text
                                break
                        if slide_title:
                            break
                
                if slide_title:
                    d.text((40, 80), f"Title: {slide_title}", fill=(0, 0, 0), font=font)
                
                # Save the image
                img.save(slide_png_path)
                png_files.append(slide_png_path)
                
                print(f"Created PNG for slide {i+1}: {slide_png_path}")
            
            print(f"Generated {len(png_files)} PNG files directly from PPTX")
            
        except Exception as direct_conversion_error:
            print(f"Error in direct PPTX to PNG conversion: {str(direct_conversion_error)}")
            import traceback
            print(traceback.format_exc())
            
            # Last resort: create a single fallback PNG
            fallback_png_path = os.path.join(output_dir, f"{basename}_fallback.png")
            
            try:
                # Create a simple image with text
                img = Image.new('RGB', (800, 600), color=(255, 255, 255))
                d = ImageDraw.Draw(img)
                d.text((10, 10), f"Failed to convert presentation. Please check the PPTX file.", fill=(0, 0, 0))
                img.save(fallback_png_path)
                
                png_files = [fallback_png_path]
                print(f"Created fallback PNG: {fallback_png_path}")
            except Exception as fallback_error:
                print(f"Error creating fallback PNG: {str(fallback_error)}")
    
    # Clean up the PDF file if it exists and we don't need it anymore
    if pdf_success and os.path.exists(pdf_path) and png_files:
        try:
            os.remove(pdf_path)
            print(f"Removed temporary PDF file: {pdf_path}")
        except Exception as e:
            print(f"Warning: Failed to remove temporary PDF file: {str(e)}")
    
    return png_files 