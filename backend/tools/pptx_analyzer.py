import os
import sys
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_font_info(paragraph):
    """
    Extract font information from a paragraph.
    """
    if not paragraph or not paragraph.runs or len(paragraph.runs) == 0:
        return "No font information available", {}
    
    run = paragraph.runs[0]  # Get the first run
    font = run.font
    
    info = []
    font_dict = {}
    
    if hasattr(font, 'name') and font.name:
        info.append(f"Font: {font.name}")
        font_dict['name'] = font.name
    if hasattr(font, 'size') and font.size:
        size_pt = font.size.pt if hasattr(font.size, 'pt') else "unknown"
        info.append(f"Size: {size_pt}pt")
        font_dict['size'] = size_pt
    if hasattr(font, 'bold') and font.bold is not None:
        info.append(f"Bold: {font.bold}")
        font_dict['bold'] = font.bold
    if hasattr(font, 'italic') and font.italic is not None:
        info.append(f"Italic: {font.italic}")
        font_dict['italic'] = font.italic
    if hasattr(font, 'underline') and font.underline is not None:
        info.append(f"Underline: {font.underline}")
        font_dict['underline'] = font.underline
    if hasattr(font, 'color') and font.color and hasattr(font.color, 'rgb'):
        rgb = font.color.rgb
        if rgb:
            hex_color = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}" if len(rgb) >= 3 else "unknown"
            info.append(f"Color: {hex_color}")
            font_dict['color'] = hex_color
    
    return ", ".join(info) if info else "No specific font properties found", font_dict

def analyze_shape_text_frames(shape, indent=""):
    """
    Analyze text frames in a shape and extract paragraph and font information.
    """
    results = []
    data = []
    
    if hasattr(shape, 'text_frame') and shape.text_frame:
        text_frame = shape.text_frame
        
        # Get paragraph information
        if hasattr(text_frame, 'paragraphs'):
            for i, paragraph in enumerate(text_frame.paragraphs):
                par_text = paragraph.text
                par_text_display = par_text
                if par_text and len(par_text) > 30:
                    par_text_display = par_text[:30] + "..."
                
                # Get paragraph level (indentation)
                level = getattr(paragraph, 'level', '0')
                
                # Get paragraph alignment
                alignment = getattr(paragraph, 'alignment', 'Unknown')
                
                # Get font information
                font_info_str, font_info_dict = get_font_info(paragraph)
                
                results.append(f"{indent}Paragraph {i+1}: '{par_text_display}', Level: {level}, Alignment: {alignment}")
                results.append(f"{indent}  Font: {font_info_str}")
                
                # Build data dictionary
                para_dict = {
                    'index': i+1,
                    'text': par_text,
                    'level': level,
                    'alignment': str(alignment),
                    'font': font_info_dict
                }
                data.append(para_dict)
    
    return results, data

def analyze_layout(prs, layout_name, export_json=False):
    """
    Analyze a specific layout by name and output detailed information.
    """
    print(f"\n===== ANALYZING LAYOUT: {layout_name} =====")
    
    # Create result data structure
    layout_data = {
        'name': layout_name,
        'placeholders': [],
        'shapes': []
    }
    
    # Find layout by name
    target_layout = None
    layout_index = None
    
    for i, layout in enumerate(prs.slide_layouts):
        curr_name = getattr(layout, 'name', f'Layout {i+1}')
        if layout_name.lower() in curr_name.lower():
            target_layout = layout
            layout_index = i
            break
    
    if not target_layout:
        print(f"ERROR: Layout '{layout_name}' not found")
        print("Available layouts:")
        layouts = []
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = getattr(layout, 'name', f'Layout {i+1}')
            print(f"  {i+1}. {layout_name}")
            layouts.append({'index': i+1, 'name': layout_name})
        
        layout_data['error'] = f"Layout '{layout_name}' not found"
        layout_data['available_layouts'] = layouts
        
        if export_json:
            return layout_data
        return
    
    print(f"Found layout '{layout_name}' at index {layout_index}")
    layout_data['index'] = layout_index
    layout_data['full_name'] = getattr(target_layout, 'name', f'Layout {layout_index+1}')
    
    # Print placeholders in this layout
    print(f"\nPlaceholders ({len(target_layout.placeholders)}):")
    for i, placeholder in enumerate(target_layout.placeholders):
        idx = placeholder.placeholder_format.idx if hasattr(placeholder, 'placeholder_format') else 'Unknown'
        ph_type = placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else 'Unknown'
        name = getattr(placeholder, 'name', f'Placeholder {i+1}')
        print(f"  {i+1}. Name: '{name}', Type: {ph_type}, Index: {idx}")
        
        placeholder_data = {
            'index': i+1,
            'name': name,
            'placeholder_type': str(ph_type),
            'placeholder_index': str(idx),
            'paragraphs': []
        }
        
        # Get text frame information
        text_info, para_data = analyze_shape_text_frames(placeholder, indent="    ")
        for line in text_info:
            print(line)
        
        placeholder_data['paragraphs'] = para_data
            
        # Detailed internal XML structure analysis
        print(f"    Internal structure:")
        xml_info = {}
        if hasattr(placeholder, 'element'):
            try:
                # Check for placeholder's inheritance from master
                print(f"    XML inheritance check:")
                if placeholder.element.getparent() is not None:
                    parent_type = placeholder.element.getparent().tag
                    print(f"      Parent element type: {parent_type}")
                    xml_info['parent_element_type'] = parent_type
                
                # Try to see if this is linked to the master slide
                is_from_master = hasattr(placeholder, 'follow_master_graphics')
                print(f"      Is from master slide: {is_from_master}")
                xml_info['is_from_master'] = is_from_master
                
                # Try to check formatting directly in XML
                has_tx_body = False
                if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                    if hasattr(placeholder.text_frame, '_txBody'):
                        has_tx_body = True
                        print(f"      Has _txBody: Yes")
                xml_info['has_tx_body'] = has_tx_body
                
                # Check for run-level properties in XML
                run_elements = []
                if hasattr(placeholder, 'text_frame') and placeholder.text_frame and placeholder.text_frame.paragraphs:
                    for j, para in enumerate(placeholder.text_frame.paragraphs):
                        if para.runs:
                            for k, run in enumerate(para.runs):
                                has_r_element = hasattr(run, '_r')
                                if has_r_element:
                                    print(f"      Paragraph {j+1}, Run {k+1} has _r element: Yes")
                                    run_elements.append({
                                        'paragraph': j+1,
                                        'run': k+1,
                                        'has_r_element': has_r_element
                                    })
                xml_info['run_elements'] = run_elements
            except Exception as e:
                print(f"      Error analyzing XML: {str(e)}")
                xml_info['error'] = str(e)
        
        placeholder_data['xml_info'] = xml_info
        layout_data['placeholders'].append(placeholder_data)
    
    # Analyze all shapes in the layout
    print(f"\nShapes ({len(target_layout.shapes)}):")
    for i, shape in enumerate(target_layout.shapes):
        shape_type = "Unknown"
        shape_type_enum = None
        if hasattr(shape, 'shape_type'):
            shape_type_enum = shape.shape_type
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                shape_type = "PLACEHOLDER"
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shape_type = "AUTO_SHAPE"
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                shape_type = "TEXT_BOX"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_type = "PICTURE"
            else:
                shape_type = str(shape.shape_type)
        
        name = getattr(shape, 'name', f'Shape {i+1}')
        text = ""
        if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
            text = shape.text_frame.text
            text_display = text
            if len(text) > 30:
                text_display = text[:30] + "..."
        else:
            text_display = ""
        
        print(f"  {i+1}. Name: '{name}', Type: {shape_type}, Text: '{text_display}'")
        
        shape_data = {
            'index': i+1,
            'name': name,
            'shape_type': shape_type,
            'shape_type_enum': str(shape_type_enum),
            'text': text,
            'paragraphs': []
        }
        
        # Get text frame information
        text_info, para_data = analyze_shape_text_frames(shape, indent="    ")
        for line in text_info:
            print(line)
        
        shape_data['paragraphs'] = para_data
            
        # Look for run-level style inheritance details
        style_info = []
        if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.paragraphs:
            for j, para in enumerate(shape.text_frame.paragraphs):
                para_style = {'paragraph': j+1, 'runs': []}
                
                if para.runs:
                    for k, run in enumerate(para.runs):
                        run_style = {'run': k+1}
                        
                        if hasattr(run, 'font'):
                            font = run.font
                            has_rPr = hasattr(font, '_rPr')
                            if has_rPr:
                                print(f"    Paragraph {j+1}, Run {k+1} has style definition (_rPr): Yes")
                                run_style['has_rPr'] = True
                            else:
                                run_style['has_rPr'] = False
                        
                        para_style['runs'].append(run_style)
                
                # Check for paragraph style inheritance
                has_pPr = False
                if hasattr(para, '_p'):
                    if hasattr(para._p, 'pPr'):
                        has_pPr = True
                        print(f"    Paragraph {j+1} has paragraph properties (pPr): Yes")
                
                para_style['has_pPr'] = has_pPr
                style_info.append(para_style)
        
        shape_data['style_info'] = style_info
        layout_data['shapes'].append(shape_data)
    
    if export_json:
        return layout_data
    
    return None

def analyze_pptx_template(template_path, export_json=False, output_path=None):
    """
    Analyze a PowerPoint template file and output detailed information about layouts and shapes.
    """
    if not os.path.exists(template_path):
        print(f"ERROR: Template file does not exist: {template_path}")
        if export_json:
            return {'error': f"Template file does not exist: {template_path}"}
        return
    
    # Load the presentation
    print(f"\n===== ANALYZING TEMPLATE: {template_path} =====")
    prs = Presentation(template_path)
    
    # Create result data structure
    result = {
        'template_path': template_path,
        'slide_masters': [],
        'layouts': [],
        'toc_layouts': []
    }
    
    # Analyze slide masters
    print(f"\n===== SLIDE MASTERS ({len(prs.slide_masters)}) =====")
    for i, master in enumerate(prs.slide_masters):
        print(f"\nMaster {i+1}:")
        master_name = getattr(master, 'name', 'Unknown')
        print(f"  Name: {master_name}")
        
        master_data = {
            'index': i+1,
            'name': master_name,
            'layouts': []
        }
        
        # Analyze slide layouts in this master
        print(f"\n  Layouts in this master ({len(master.slide_layouts)}):")
        for j, layout in enumerate(master.slide_layouts):
            layout_name = getattr(layout, 'name', f'Layout {j+1}')
            print(f"\n    Layout {j+1}: {layout_name}")
            
            layout_data = {
                'index': j+1,
                'name': layout_name,
                'placeholders': [],
                'shapes': []
            }
            
            # Print placeholders in this layout
            print(f"      Placeholders ({len(layout.placeholders)}):")
            for k, placeholder in enumerate(layout.placeholders):
                idx = placeholder.placeholder_format.idx if hasattr(placeholder, 'placeholder_format') else 'Unknown'
                ph_type = placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else 'Unknown'
                name = getattr(placeholder, 'name', f'Placeholder {k+1}')
                print(f"        {k+1}. Name: '{name}', Type: {ph_type}, Index: {idx}")
                
                placeholder_data = {
                    'index': k+1,
                    'name': name,
                    'placeholder_type': str(ph_type),
                    'placeholder_index': str(idx),
                    'paragraphs': []
                }
                
                # Get text frame information
                text_info, para_data = analyze_shape_text_frames(placeholder, indent="          ")
                for line in text_info:
                    print(line)
                
                placeholder_data['paragraphs'] = para_data
                layout_data['placeholders'].append(placeholder_data)
            
            # Analyze all shapes in the layout
            print(f"      Shapes ({len(layout.shapes)}):")
            for k, shape in enumerate(layout.shapes):
                shape_type = "Unknown"
                shape_type_enum = None
                if hasattr(shape, 'shape_type'):
                    shape_type_enum = shape.shape_type
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        shape_type = "PLACEHOLDER"
                    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        shape_type = "AUTO_SHAPE"
                    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        shape_type = "TEXT_BOX"
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        shape_type = "PICTURE"
                    else:
                        shape_type = str(shape.shape_type)
                
                name = getattr(shape, 'name', f'Shape {k+1}')
                text = ""
                if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                    text = shape.text_frame.text
                    text_display = text
                    if len(text) > 30:
                        text_display = text[:30] + "..."
                else:
                    text_display = ""
                
                print(f"        {k+1}. Name: '{name}', Type: {shape_type}, Text: '{text_display}'")
                
                shape_data = {
                    'index': k+1,
                    'name': name,
                    'shape_type': shape_type,
                    'shape_type_enum': str(shape_type_enum),
                    'text': text,
                    'paragraphs': []
                }
                
                # Get text frame information
                text_info, para_data = analyze_shape_text_frames(shape, indent="          ")
                for line in text_info:
                    print(line)
                
                shape_data['paragraphs'] = para_data
                layout_data['shapes'].append(shape_data)
            
            master_data['layouts'].append(layout_data)
        
        result['slide_masters'].append(master_data)
    
    # Analyze the default layouts (shortcut to slide_masters[0].slide_layouts)
    print(f"\n===== DEFAULT LAYOUTS ({len(prs.slide_layouts)}) =====")
    for i, layout in enumerate(prs.slide_layouts):
        layout_name = getattr(layout, 'name', f'Layout {i+1}')
        print(f"\nLayout {i+1}: '{layout_name}'")
        
        layout_data = {
            'index': i+1,
            'name': layout_name,
            'placeholders': []
        }
        
        # Print placeholders in this layout
        print(f"  Placeholders ({len(layout.placeholders)}):")
        for j, placeholder in enumerate(layout.placeholders):
            idx = placeholder.placeholder_format.idx if hasattr(placeholder, 'placeholder_format') else 'Unknown'
            ph_type = placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else 'Unknown'
            name = getattr(placeholder, 'name', f'Placeholder {j+1}')
            print(f"    {j+1}. Name: '{name}', Type: {ph_type}, Index: {idx}")
            
            placeholder_data = {
                'index': j+1,
                'name': name,
                'placeholder_type': str(ph_type),
                'placeholder_index': str(idx),
                'paragraphs': []
            }
            
            # Get text frame information
            text_info, para_data = analyze_shape_text_frames(placeholder, indent="      ")
            for line in text_info:
                print(line)
            
            placeholder_data['paragraphs'] = para_data
            layout_data['placeholders'].append(placeholder_data)
        
        result['layouts'].append(layout_data)

    # Find Table of Contents related layouts
    print("\n===== TABLE OF CONTENTS SPECIFIC ANALYSIS =====")
    toc_layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layout_name = getattr(layout, 'name', '').lower()
        if 'content' in layout_name or 'toc' in layout_name or 'table' in layout_name:
            toc_layouts.append((i, layout))
    
    if not toc_layouts:
        print("No layouts with names related to Table of Contents found!")
        print("Looking for layouts with section-related shapes instead...")
        
        # Look for layouts with section shapes
        for i, layout in enumerate(prs.slide_layouts):
            section_related = False
            for shape in layout.shapes:
                name = getattr(shape, 'name', '').lower()
                if 'section' in name or ('0' in name and len(name) <= 3):
                    section_related = True
                    break
            
            if section_related:
                toc_layouts.append((i, layout))
    
    if not toc_layouts:
        print("Could not identify any Table of Contents related layouts!")
        result['toc_layouts'] = []
    else:
        print(f"Found {len(toc_layouts)} potential Table of Contents layouts:")
        for i, (idx, layout) in enumerate(toc_layouts):
            layout_name = getattr(layout, 'name', f'Layout {idx+1}')
            print(f"\nPotential TOC Layout {i+1}: '{layout_name}' (Index: {idx})")
            
            toc_layout_data = {
                'index': i+1,
                'layout_index': idx,
                'name': layout_name,
                'section_shapes': []
            }
            
            # Look for section shapes
            section_shapes = []
            for j, shape in enumerate(layout.shapes):
                name = getattr(shape, 'name', '').lower()
                if ('section' in name or 
                    ('0' in name and len(name) <= 3) or 
                    'divider' in name):
                    shape_type = "Unknown"
                    if hasattr(shape, 'shape_type'):
                        shape_type = str(shape.shape_type)
                    
                    text = ""
                    if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                        text = shape.text_frame.text
                        text_display = text
                        if len(text) > 30:
                            text_display = text[:30] + "..."
                    else:
                        text_display = ""
                    
                    section_shapes.append((j, shape, name, shape_type, text))
                    
                    section_shape_data = {
                        'index': j+1,
                        'name': name,
                        'shape_type': shape_type,
                        'text': text,
                        'paragraphs': []
                    }
                    
                    # Get text frame information
                    text_info, para_data = analyze_shape_text_frames(shape, indent="      ")
                    for line in text_info:
                        print(line)
                    
                    section_shape_data['paragraphs'] = para_data
                    toc_layout_data['section_shapes'].append(section_shape_data)
            
            if section_shapes:
                print(f"  Found {len(section_shapes)} section-related shapes:")
                for j, shape, name, shape_type, text in section_shapes:
                    print(f"    {j+1}. Name: '{name}', Type: {shape_type}, Text: '{text_display}'")
                
                # Analyze section naming patterns
                print("\n  Section naming patterns analysis:")
                section_name_patterns = {}
                for _, _, name, _, _ in section_shapes:
                    # Try to extract patterns
                    if 'section' in name:
                        # Check if it follows pattern like "section1", "Section2", etc.
                        import re
                        match = re.search(r'section(\d+)', name, re.IGNORECASE)
                        if match:
                            pattern = name.replace(match.group(1), 'X')
                            if pattern not in section_name_patterns:
                                section_name_patterns[pattern] = []
                            section_name_patterns[pattern].append(name)
                    elif name.isdigit() or (len(name) <= 3 and '0' in name):
                        # Check if it's just a number like "01", "02"
                        pattern = 'XX'
                        if pattern not in section_name_patterns:
                            section_name_patterns[pattern] = []
                        section_name_patterns[pattern].append(name)
                    elif 'divider' in name:
                        # Check if it follows pattern like "01Divider", "divider2"
                        import re
                        match = re.search(r'(\d+)divider', name, re.IGNORECASE)
                        if match:
                            pattern = name.replace(match.group(1), 'X')
                            if pattern not in section_name_patterns:
                                section_name_patterns[pattern] = []
                            section_name_patterns[pattern].append(name)
                
                pattern_data = {}
                if section_name_patterns:
                    print("    Detected section naming patterns:")
                    for pattern, examples in section_name_patterns.items():
                        print(f"    - Pattern: '{pattern}', Examples: {examples[:3]}")
                        pattern_data[pattern] = examples
                    
                    # Create suggestions for patterns
                    suggestions = {
                        'section_patterns': [],
                        'number_patterns': [],
                        'divider_patterns': []
                    }
                    
                    print("\n  Suggested get_toc_shapes patterns to use in generate_pptx.py:")
                    section_patterns = [p for p in section_name_patterns if 'section' in p.lower()]
                    number_patterns = [p for p in section_name_patterns if p == 'XX']
                    divider_patterns = [p for p in section_name_patterns if 'divider' in p.lower()]
                    
                    if section_patterns:
                        print("    Section name patterns:")
                        for pattern in section_patterns:
                            suggestion = f"{pattern.replace('X', '{section_num_str}')}"
                            print(f"        f\"{suggestion}\"")
                            suggestions['section_patterns'].append(suggestion)
                    
                    if number_patterns:
                        print("    Number name patterns:")
                        suggestion1 = "{section_num_str}"
                        print(f"        f\"{suggestion1}\"")
                        suggestions['number_patterns'].append(suggestion1)
                        
                        if len(section_name_patterns['XX'][0]) == 2:
                            suggestion2 = "{section_num_str.zfill(2)}"
                            print(f"        f\"{suggestion2}\"")
                            suggestions['number_patterns'].append(suggestion2)
                    
                    if divider_patterns:
                        print("    Divider name patterns:")
                        for pattern in divider_patterns:
                            suggestion = f"{pattern.replace('X', '{section_num_str}')}"
                            print(f"        f\"{suggestion}\"")
                            suggestions['divider_patterns'].append(suggestion)
                
                toc_layout_data['patterns'] = pattern_data
                toc_layout_data['suggestions'] = suggestions if section_name_patterns else {}
            else:
                print("  No section-related shapes found in this layout.")
            
            result['toc_layouts'].append(toc_layout_data)
    
    print("\n===== ANALYSIS COMPLETE =====")
    
    # Export to JSON if requested
    if export_json:
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(result, f, indent=2)
            print(f"\nExported analysis to {output_path}")
        return result
    
    return None

if __name__ == "__main__":
    # Parse command line arguments
    import argparse
    parser = argparse.ArgumentParser(description='Analyze PowerPoint templates')
    parser.add_argument('template_path', help='Path to the PowerPoint template file')
    parser.add_argument('-l', '--layout', help='Name of specific layout to analyze')
    parser.add_argument('-j', '--json', action='store_true', help='Export analysis as JSON')
    parser.add_argument('-o', '--output', help='Output file path for JSON export')
    
    args = parser.parse_args()
    
    if args.layout:
        # Analyze specific layout
        prs = Presentation(args.template_path)
        result = analyze_layout(prs, args.layout, export_json=args.json)
        
        if args.json and result:
            output_path = args.output or f"layout_{args.layout.replace(' ', '_').lower()}_analysis.json"
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(result, f, indent=2)
            print(f"\nExported layout analysis to {output_path}")
    else:
        # Analyze entire template
        output_path = args.output or f"{os.path.splitext(os.path.basename(args.template_path))[0]}_analysis.json"
        analyze_pptx_template(args.template_path, export_json=args.json, output_path=output_path if args.json else None) 