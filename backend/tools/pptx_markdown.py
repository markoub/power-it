"""
Functions for handling Markdown formatting in PowerPoint slides
"""
import re
import markdown
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

def parse_markdown_to_runs(paragraph, markdown_text):
    """
    Parse markdown text and add appropriate runs to a paragraph with formatting.
    
    Args:
        paragraph: The PowerPoint paragraph object to add runs to
        markdown_text: The markdown text to parse
        
    Returns:
        None (modifies paragraph in place)
    """
    # Clear any existing text
    if paragraph.runs:
        for _ in range(len(paragraph.runs)):
            paragraph.runs[0].text = ""
    
    # Handle basic markdown with regex
    # This is a simplified implementation - more complex markdown would require a proper parser
    text_segments = []
    
    # Process bold text: **bold** or __bold__
    pattern = r'\*\*(.*?)\*\*|__(.*?)__'
    last_end = 0
    for match in re.finditer(pattern, markdown_text):
        # Add any text before this match
        if match.start() > last_end:
            text_segments.append({
                'text': markdown_text[last_end:match.start()],
                'bold': False,
                'italic': False
            })
        
        # Add the bold text (using either group that matched)
        bold_text = match.group(1) if match.group(1) is not None else match.group(2)
        text_segments.append({
            'text': bold_text,
            'bold': True,
            'italic': False
        })
        
        last_end = match.end()
    
    # Add any remaining text after the last bold match
    if last_end < len(markdown_text):
        text_segments.append({
            'text': markdown_text[last_end:],
            'bold': False,
            'italic': False
        })
    
    # Now process italic text in each segment
    processed_segments = []
    for segment in text_segments:
        if segment['bold']:
            # Keep bold segments as they are
            processed_segments.append(segment)
            continue
            
        text = segment['text']
        last_end = 0
        italic_pattern = r'\*(.*?)\*|_(.*?)_'
        
        italic_found = False
        for match in re.finditer(italic_pattern, text):
            italic_found = True
            # Add any text before this match
            if match.start() > last_end:
                processed_segments.append({
                    'text': text[last_end:match.start()],
                    'bold': False,
                    'italic': False
                })
            
            # Add the italic text (using either group that matched)
            italic_text = match.group(1) if match.group(1) is not None else match.group(2)
            processed_segments.append({
                'text': italic_text,
                'bold': False,
                'italic': True
            })
            
            last_end = match.end()
        
        # Add any remaining text
        if not italic_found:
            processed_segments.append(segment)
        elif last_end < len(text):
            processed_segments.append({
                'text': text[last_end:],
                'bold': False,
                'italic': False
            })
    
    # Add all processed segments as runs
    for segment in processed_segments:
        if not segment['text']:  # Skip empty segments
            continue
            
        run = paragraph.add_run()
        run.text = segment['text']
        
        if segment['bold']:
            run.font.bold = True
        
        if segment['italic']:
            run.font.italic = True
    
    print(f"Processed Markdown text with {len(processed_segments)} formatted segments")

def apply_markdown_to_text_frame(text_frame, markdown_text, base_size=18):
    """
    Apply markdown formatting to a text frame.
    
    Args:
        text_frame: The PowerPoint text frame object
        markdown_text: The markdown text to apply
        base_size: Base font size for text
        
    Returns:
        None (modifies text_frame in place)
    """
    # Clear existing paragraphs except for the first one
    while len(text_frame.paragraphs) > 1:
        p = text_frame.paragraphs[-1]
        tr = p._element
        tr.getparent().remove(tr)
    
    # Use the first paragraph
    paragraph = text_frame.paragraphs[0]
    
    # Clear existing text/runs
    if paragraph.runs:
        for _ in range(len(paragraph.runs)):
            paragraph.runs[0].text = ""
            
    # Set default paragraph properties
    paragraph.font.size = Pt(base_size)
    
    # Apply the markdown formatting
    parse_markdown_to_runs(paragraph, markdown_text)

def format_markdown_content(content_shape, content_items):
    """Format content text with markdown processing and appropriate styles"""
    if not content_shape or not hasattr(content_shape, "text_frame"):
        return
    
    # Clear the content placeholder
    content_shape.text_frame.clear()
    
    # Apply word wrapping
    content_shape.text_frame.word_wrap = True
    
    # Determine appropriate font size based on total content length and item count
    total_length = sum(len(item) for item in content_items)
    num_items = len(content_items)
    
    print(f"Formatting markdown content with {num_items} items, total length: {total_length}")
    
    # Determine base font size based on content volume
    base_size = 18  # Default size
    if total_length > 800 or num_items > 10:
        base_size = 12  # Very dense content
        print(f"  Using small font (12pt) for very dense content")
    elif total_length > 500 or num_items > 7:
        base_size = 14  # Dense content
        print(f"  Using reduced font (14pt) for dense content")
    elif total_length > 300 or num_items > 5:
        base_size = 16  # Moderate content
        print(f"  Using medium font (16pt) for moderate content")
    else:
        print(f"  Using standard font (18pt) for normal content")
    
    # Add each content item as a paragraph with markdown formatting
    for idx, item in enumerate(content_items):
        if idx == 0:
            p = content_shape.text_frame.paragraphs[0]
        else:
            p = content_shape.text_frame.add_paragraph()
        
        # Set base font size for paragraph
        p.font.size = Pt(base_size)
        p.level = 0
        
        # Add a bit of spacing between items but not too much
        if idx > 0:
            p.space_before = Pt(6)
        
        # Apply markdown formatting to this paragraph
        parse_markdown_to_runs(p, item)
    
    print(f"  Added {len(content_items)} formatted markdown content items") 