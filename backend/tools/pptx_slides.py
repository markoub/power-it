"""
Functions for handling different PowerPoint slide types
"""
import os
from typing import Dict, Any, Optional
from pptx.util import Pt, Inches

from .pptx_shapes import find_shape_by_name
from .pptx_text import adjust_text_size
from .pptx_utils import format_section_number
from .pptx_markdown import parse_markdown_to_runs, format_markdown_content

def create_welcome_slide(prs, welcome_layout, presentation_title, welcome_data=None):
    """
    Create a welcome/title slide for the presentation.
    
    Args:
        prs: The presentation object
        welcome_layout: The slide layout to use
        presentation_title: The presentation title
        welcome_data: Optional dictionary with title, subtitle, and author fields
    
    Returns:
        The created welcome slide
    """
    print("\n===== Creating Welcome slide =====")
    welcome_slide = prs.slides.add_slide(welcome_layout)
    
    # Find title placeholder
    title_placeholder = None
    subtitle_placeholder = None
    author_placeholder = None
    
    print("Locating welcome slide placeholders:")
    for i, shape in enumerate(welcome_slide.shapes):
        shape_name = getattr(shape, "name", "Unknown")
        shape_type = "Unknown"
        if hasattr(shape, "placeholder_format") and hasattr(shape.placeholder_format, "type"):
            shape_type = shape.placeholder_format.type
            
        shape_text = ""
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            shape_text = shape.text_frame.text
            
        print(f"  Shape {i}: Name='{shape_name}', Type={shape_type}, Text='{shape_text}'")
        
        # Check for title placeholder
        if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # 1 is TITLE
            title_placeholder = shape
            print(f"  Found welcome slide title placeholder by type, name: {shape_name}")
        elif hasattr(shape, "name") and shape.name and "title" in shape.name.lower():
            title_placeholder = shape
            print(f"  Found welcome slide title placeholder by name: {shape_name}")
            
        # Check for subtitle placeholder
        if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 2:  # 2 is SUBTITLE
            subtitle_placeholder = shape
            print(f"  Found welcome slide subtitle placeholder by type, name: {shape_name}")
        elif hasattr(shape, "name") and shape.name and "subtitle" in shape.name.lower():
            subtitle_placeholder = shape
            print(f"  Found welcome slide subtitle placeholder by name: {shape_name}")
            
        # Check for author placeholder
        if hasattr(shape, "name") and shape.name and "author" in shape.name.lower():
            author_placeholder = shape
            print(f"  Found welcome slide author placeholder by name: {shape_name}")
    
    # Set the presentation title on the welcome slide
    if title_placeholder and hasattr(title_placeholder, "text_frame"):
        # Use welcome_data title if available, otherwise use presentation title
        title_text = welcome_data.get("title", presentation_title) if welcome_data else presentation_title
        
        # Clear the title placeholder
        title_placeholder.text_frame.clear()
        p = title_placeholder.text_frame.paragraphs[0]
        
        # Apply markdown formatting to title
        parse_markdown_to_runs(p, title_text)
        print(f"Set welcome slide title to: {title_text}")
        
        # Adjust title size based on length
        adjust_text_size(title_placeholder.text_frame, default_size=40, min_size=28, max_length=40, long_text_size=32, is_title=True)
    else:
        print("⚠️ WARNING: Could not find Title placeholder on welcome slide.")
    
    # Set subtitle if available
    if subtitle_placeholder and hasattr(subtitle_placeholder, "text_frame") and welcome_data and "subtitle" in welcome_data:
        # Clear the subtitle placeholder
        subtitle_placeholder.text_frame.clear()
        p = subtitle_placeholder.text_frame.paragraphs[0]
        
        # Apply markdown formatting to subtitle
        parse_markdown_to_runs(p, welcome_data["subtitle"])
        print(f"Set welcome slide subtitle to: {welcome_data['subtitle']}")
        
        # Adjust subtitle size based on length
        adjust_text_size(subtitle_placeholder.text_frame, default_size=24, min_size=18, max_length=60, long_text_size=20, is_title=True)
    
    # Set author if available
    if author_placeholder and hasattr(author_placeholder, "text_frame") and welcome_data and "author" in welcome_data:
        author_placeholder.text_frame.text = welcome_data["author"]
        print(f"Set welcome slide author to: {welcome_data['author']}")
        
    return welcome_slide

def create_section_slide(prs, section_layout, section_title, section_number):
    """
    Create a section slide with title and section number.
    
    Args:
        prs: The presentation object
        section_layout: The slide layout to use
        section_title: The section title
        section_number: The section number (1-based)
        
    Returns:
        The created section slide
    """
    print(f"\n===== Creating Section slide {section_number}: {section_title} =====")
    slide_obj = prs.slides.add_slide(section_layout)
    
    # Print all shapes in the slide for debugging
    print("Available shapes in section slide:")
    for i, shape in enumerate(slide_obj.shapes):
        shape_name = getattr(shape, "name", f"Shape {i}")
        shape_type = getattr(shape, "shape_type", "Unknown")
        placeholder_type = "Unknown"
        if hasattr(shape, "placeholder_format"):
            placeholder_type = getattr(shape.placeholder_format, "type", "Unknown")
        shape_text = ""
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            shape_text = shape.text_frame.text
        print(f"  Shape {i}: Name='{shape_name}', Type={shape_type}, Placeholder Type={placeholder_type}, Text='{shape_text}'")
    
    # Based on layout_analysis.json, for Section layout:
    # - "Title" is placeholder type BODY (2) and should contain the section title
    # - "Number" is placeholder type CENTER_TITLE (3) and should contain the section number
    
    # Find shapes by placeholder type which is more reliable than name
    title_placeholder = None  # For section title (BODY type)
    number_placeholder = None  # For section number (CENTER_TITLE type)
    
    # First, look for placeholders by type
    for shape in slide_obj.shapes:
        if not hasattr(shape, "placeholder_format") or not hasattr(shape, "text_frame"):
            continue
            
        placeholder_type = shape.placeholder_format.type
        if placeholder_type == 2:  # BODY for title
            title_placeholder = shape
            print(f"Found title placeholder (BODY) by type: {getattr(shape, 'name', 'Unknown')}")
        elif placeholder_type == 3:  # CENTER_TITLE for number
            number_placeholder = shape
            print(f"Found number placeholder (CENTER_TITLE) by type: {getattr(shape, 'name', 'Unknown')}")
    
    # If we couldn't find by type, try to find by name with find_shape_by_name
    if not title_placeholder:
        title_placeholder = find_shape_by_name(slide_obj, "Title")
        
    if not number_placeholder:
        number_placeholder = find_shape_by_name(slide_obj, "Number")
    
    # Last resort: just get the first two shapes with text frames
    if not title_placeholder and not number_placeholder and len(slide_obj.shapes) >= 2:
        text_shapes = [s for s in slide_obj.shapes if hasattr(s, "text_frame")]
        if len(text_shapes) >= 2:
            title_placeholder = text_shapes[0]  # Use first shape for title
            number_placeholder = text_shapes[1]  # Use second shape for number
            print("Using first two text shapes as fallback placeholders")
    
    # Set section title
    if title_placeholder and hasattr(title_placeholder, "text_frame"):
        # First ensure any existing content is cleared
        title_placeholder.text_frame.clear()
        
        # Then add a paragraph with the section title
        p = title_placeholder.text_frame.paragraphs[0]
        
        # Apply markdown formatting to section title
        parse_markdown_to_runs(p, section_title)
        print(f"  Set section title to: {section_title}")
        
        # Adjust section title size based on length
        adjust_text_size(title_placeholder.text_frame, default_size=32, min_size=24, max_length=35, long_text_size=28, is_title=True)
    else:
        print("⚠️ ERROR: Could not find a valid title placeholder on section slide")
    
    # Set section number
    if number_placeholder and hasattr(number_placeholder, "text_frame"):
        # First ensure any existing content is cleared
        number_placeholder.text_frame.clear()
        
        # Then add a paragraph with the section number
        p = number_placeholder.text_frame.paragraphs[0]
        p.text = format_section_number(section_number)
        print(f"  Set section number to: {format_section_number(section_number)}")
    else:
        print("⚠️ ERROR: Could not find a valid number placeholder on section slide")
        
    return slide_obj

def create_content_slide(prs, content_layout, slide_title, content, add_image=False, image_path=None):
    """
    Create a content slide with title and bullet points.
    
    Args:
        prs: The presentation object
        content_layout: The slide layout to use
        slide_title: The slide title
        content: List of content items (bullet points)
        add_image: Whether to add an image to the slide
        image_path: Path to the image file
        
    Returns:
        The created content slide
    """
    print(f"\n===== Creating Content slide: {slide_title} =====")
    slide_obj = prs.slides.add_slide(content_layout)
    
    # Find title shape
    title_shape = None
    for shape in slide_obj.shapes:
        if shape.name.lower().startswith('title'):
            title_shape = shape
            break
            
    if not title_shape:
        for shape in slide_obj.placeholders:
            if shape.placeholder_format.type == 1:  # TITLE
                title_shape = shape
                break
    
    # Set title if found
    if title_shape:
        title_shape.text = slide_title
    else:
        print("  Warning: No title shape found in slide layout")
    
    # Find body/content shape for bullet points
    content_shape = None
    for shape in slide_obj.shapes:
        if shape.name.lower().startswith('content') or shape.name.lower().startswith('body'):
            content_shape = shape
            break
    
    if not content_shape:
        for shape in slide_obj.placeholders:
            if shape.placeholder_format.type == 2:  # BODY
                content_shape = shape
                break
    
    # Set content bullet points
    if content_shape and content:
        text_frame = content_shape.text_frame
        text_frame.clear()  # Clear existing text
        
        # Instead of using format_content_text, manually add the content
        for i, item in enumerate(content):
            if i == 0:
                # Use the first paragraph
                p = text_frame.paragraphs[0]
            else:
                # Add a new paragraph for each subsequent item
                p = text_frame.add_paragraph()
            
            # Set the paragraph text
            p.text = item
            # Set bullet level to 0 (top level)
            p.level = 0
            # Ensure font size is appropriate
            if p.runs:
                p.runs[0].font.size = Pt(18)  # Default size
                
        print(f"  Added {len(content)} content items to slide")
    else:
        print("  Warning: No content shape found in slide layout or no content provided")
    
    # Add image if requested
    if add_image and image_path:
        print(f"  Adding image to slide: {image_path}")
        
        # Validate image path
        if not os.path.exists(image_path):
            print(f"  ERROR: Image file does not exist: {image_path}")
            return slide_obj
            
        # Check if file is empty
        if os.path.getsize(image_path) == 0:
            print(f"  ERROR: Image file is empty: {image_path}")
            return slide_obj
            
        try:
            # Try to find a picture placeholder
            picture_placeholder = None
            for shape in slide_obj.placeholders:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 18:  # PICTURE
                    picture_placeholder = shape
                    print(f"  Found picture placeholder: {shape.name}")
                    break
            
            if picture_placeholder:
                # Use the placeholder's coordinates
                left, top, width, height = picture_placeholder.left, picture_placeholder.top, picture_placeholder.width, picture_placeholder.height
                
                # Insert picture into placeholder
                try:
                    picture = slide_obj.shapes.add_picture(image_path, left, top, width, height)
                    print(f"  Image added successfully using picture placeholder")
                except Exception as e:
                    print(f"  Error adding image to placeholder: {str(e)}")
                    # Fall back to manual placement
                    picture_placeholder = None
            
            # If no picture placeholder found, add the image directly to the slide
            if not picture_placeholder:
                # Check slide size to get dimensions
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Define image position (right half of slide)
                # Typically we want images to take up about 1/3 of the slide width
                # and be positioned in the right part of the slide
                img_width = slide_width / 2.5  # 40% of slide width
                
                # Placed on right side, with some margins
                left = slide_width - img_width - Inches(0.5)  # 0.5 inches from right edge
                top = Inches(1.5)  # 1.5 inches from top
                
                # Add the picture to the slide
                try:
                    # First try to import from PIL to verify the image
                    try:
                        from PIL import Image
                        img = Image.open(image_path)
                        img_format = img.format
                        img_size = img.size
                        print(f"  Verified image: format={img_format}, size={img_size}")
                        img.close()
                    except Exception as pil_err:
                        print(f"  Warning: Failed to verify image with PIL: {str(pil_err)}")
                    
                    # Now add the image to the slide
                    picture = slide_obj.shapes.add_picture(image_path, left, top, width=img_width)
                    print(f"  Image added successfully using manual placement")
                except Exception as e:
                    print(f"  Error adding image: {str(e)}")
                    import traceback
                    print(f"  Error details: {traceback.format_exc()}")
        except Exception as outer_e:
            print(f"  Error processing image: {str(outer_e)}")
            import traceback
            print(f"  Error details: {traceback.format_exc()}")
    
    return slide_obj

def create_thank_you_slide(prs, thankyou_layout):
    """
    Create a thank you slide at the end of the presentation.
    
    Args:
        prs: The presentation object
        thankyou_layout: The slide layout to use
        
    Returns:
        The created thank you slide
    """
    print("\n===== Creating ThankYou slide =====")
    thankyou_slide = prs.slides.add_slide(thankyou_layout)
    print(f"Added ThankYou slide using layout: {getattr(thankyou_layout, 'name', 'Unknown')}")
    
    # Find title placeholder if it exists
    title_shape = None
    for shape in thankyou_slide.shapes:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # 1 is TITLE
            title_shape = shape
            break
        if hasattr(shape, "name") and shape.name and ("title" in shape.name.lower() or "thank" in shape.name.lower()):
            title_shape = shape
            break
    
    # Set thank you text if found
    if title_shape and hasattr(title_shape, "text_frame"):
        title_shape.text_frame.text = "Thank You"
        print("Set ThankYou slide title to: Thank You")
        
    return thankyou_slide

def create_3images_slide(prs, slide_layout, slide_title, image_paths, subtitles):
    """
    Create a slide with three images in a grid layout, each with its own subtitle.
    
    Args:
        prs: The presentation object
        slide_layout: The slide layout to use
        slide_title: The slide title
        image_paths: Dictionary with image1, image2, image3 paths
        subtitles: Dictionary with subtitleimage1, subtitleimage2, subtitleimage3 texts
        
    Returns:
        The created 3images slide
    """
    print(f"\n===== Creating 3Images slide: {slide_title} =====")
    slide_obj = prs.slides.add_slide(slide_layout)
    
    # Find title shape
    title_shape = None
    for shape in slide_obj.shapes:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # 1 is TITLE
            title_shape = shape
            break
        if hasattr(shape, "name") and shape.name and "title" in shape.name.lower():
            title_shape = shape
            break
    
    # Set the title
    if title_shape and hasattr(title_shape, "text_frame"):
        # First clear existing text/runs
        title_shape.text_frame.clear()
        p = title_shape.text_frame.paragraphs[0]
        
        # Apply markdown formatting to title
        parse_markdown_to_runs(p, slide_title)
        print(f"  Set slide title to: {slide_title}")
        
        # Adjust slide title size based on length
        adjust_text_size(title_shape.text_frame, default_size=28, min_size=20, max_length=40, long_text_size=24, is_title=True)
    
    # Print all shapes for debugging
    print("Shapes in 3Images slide layout:")
    for i, shape in enumerate(slide_obj.shapes):
        shape_name = getattr(shape, "name", f"Shape {i}")
        shape_type = getattr(shape, "shape_type", "Unknown")
        placeholder_type = "Unknown"
        if hasattr(shape, "placeholder_format"):
            placeholder_type = getattr(shape.placeholder_format, "type", "Unknown")
        shape_text = ""
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            shape_text = shape.text_frame.text
        print(f"  Shape {i}: Name='{shape_name}', Type={shape_type}, Placeholder Type={placeholder_type}, Text='{shape_text}'")
    
    # Find all placeholders, grouped by type
    picture_placeholders = []
    text_placeholders = []
    
    # First, collect all picture and text placeholders
    for shape in slide_obj.shapes:
        # Skip the title placeholder
        if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:
            continue
            
        # Find picture placeholders
        if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 18:
            picture_placeholders.append(shape)
            print(f"  Found picture placeholder: {getattr(shape, 'name', 'Unknown')}")
            
        # Find text placeholders for subtitles (but not the title)
        elif hasattr(shape, "text_frame"):
            if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 2:
                text_placeholders.append(shape)
                print(f"  Found text placeholder: {getattr(shape, 'name', 'Unknown')}")
            elif hasattr(shape, "name") and shape.name and any(x in shape.name.lower() for x in ["text", "subtitle", "caption"]):
                text_placeholders.append(shape)
                print(f"  Found named text placeholder: {shape.name}")
    
    # Use the first 3 picture placeholders as image placeholders
    image_placeholders = []
    if picture_placeholders:
        # Sort by position (left to right, top to bottom)
        picture_placeholders.sort(key=lambda s: (s.top, s.left))
        # Take up to 3
        image_placeholders = picture_placeholders[:3]
        print(f"  Using {len(image_placeholders)} picture placeholders")
    
    # Use the first 3 text placeholders as subtitle placeholders
    subtitle_placeholders = []
    if text_placeholders:
        # Sort by position (left to right, top to bottom)
        text_placeholders.sort(key=lambda s: (s.top, s.left))
        # Take up to 3
        subtitle_placeholders = text_placeholders[:3]
        print(f"  Using {len(subtitle_placeholders)} text placeholders")
        
    # If we couldn't find enough placeholders, fall back to looking at all shapes
    if len(image_placeholders) < 3 or len(subtitle_placeholders) < 3:
        print("  Not enough placeholders found, looking at all shapes")
        
        # Find shapes that might be picture placeholders based on size and position
        all_shapes = list(slide_obj.shapes)
        # Skip the title and already found placeholders
        remaining_shapes = [s for s in all_shapes if s not in image_placeholders + subtitle_placeholders + [title_shape]]
        
        # If we need more image placeholders
        if len(image_placeholders) < 3:
            # Find shapes that look like they might be for images (typically larger than text)
            for shape in remaining_shapes:
                if hasattr(shape, "width") and hasattr(shape, "height"):
                    # If the shape is bigger than a certain size, treat it as an image placeholder
                    if shape.width > Inches(1) and shape.height > Inches(1):
                        image_placeholders.append(shape)
                        remaining_shapes.remove(shape)
                        if len(image_placeholders) >= 3:
                            break
        
        # If we need more subtitle placeholders
        if len(subtitle_placeholders) < 3:
            # Use remaining shapes with text frames
            for shape in remaining_shapes:
                if hasattr(shape, "text_frame"):
                    subtitle_placeholders.append(shape)
                    if len(subtitle_placeholders) >= 3:
                        break
    
    # Sort placeholders by position (left to right, top to bottom) if needed
    if len(image_placeholders) >= 3:
        # Sort by position (top then left) to match image1, image2, image3 order
        image_placeholders = sorted(image_placeholders[:3], key=lambda shape: (shape.top, shape.left))
        
    if len(subtitle_placeholders) >= 3:
        # Sort by position to match subtitle1, subtitle2, subtitle3 order
        subtitle_placeholders = sorted(subtitle_placeholders[:3], key=lambda shape: (shape.top, shape.left))
    
    # Add each image if available
    for i, placeholder_name in enumerate(['image1', 'image2', 'image3']):
        if i >= len(image_placeholders):
            print(f"  No placeholder available for {placeholder_name}")
            continue
            
        # Get image path
        image_path = None
        if placeholder_name in image_paths and image_paths[placeholder_name]:
            image_path = image_paths[placeholder_name]
            if not os.path.exists(image_path):
                print(f"  Image file not found for {placeholder_name}: {image_path}")
                image_path = None
            else:
                print(f"  Image file exists for {placeholder_name}: {image_path}")
                print(f"  Image file size: {os.path.getsize(image_path)} bytes")
        
        if not image_path:
            print(f"  No valid image path for {placeholder_name}")
            continue
        
        try:
            # Get placeholder properties
            shape = image_placeholders[i]
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            
            print(f"  Placeholder dimensions for {placeholder_name} - Left: {left}, Top: {top}, Width: {width}, Height: {height}")
            
            # Try to validate the image file
            try:
                from PIL import Image
                img = Image.open(image_path)
                print(f"  Successfully opened image with PIL - Format: {img.format}, Size: {img.size}")
                img.close()
            except Exception as e:
                print(f"  Warning - Failed to validate image with PIL: {str(e)}")
            
            # First try using insert_picture if available
            if hasattr(shape, "insert_picture") and callable(shape.insert_picture):
                try:
                    shape.insert_picture(image_path)
                    print(f"  Added {placeholder_name} using insert_picture: {image_path}")
                    continue
                except Exception as e:
                    print(f"  Error with insert_picture for {placeholder_name}, trying alternative method: {str(e)}")
            
            # If insert_picture failed or isn't available, use the replacement method
            try:
                # Try to remove the placeholder before adding the image
                shape._element.getparent().remove(shape._element)
                print(f"  Successfully removed placeholder for {placeholder_name}")
            except Exception as e:
                print(f"  Error removing placeholder for {placeholder_name} (continuing anyway): {str(e)}")
            
            # Add the actual image
            added_picture = slide_obj.shapes.add_picture(image_path, left, top, width, height)
            print(f"  Added {placeholder_name} using add_picture: {image_path}")
            print(f"  Added picture object ID: {id(added_picture)}")
            
        except Exception as e:
            print(f"  Error processing {placeholder_name}: {str(e)}")
            import traceback
            print(f"  Error details:\n{traceback.format_exc()}")
    
    # Add each subtitle if available
    for i in range(3):
        # Get the correct index for subtitle placeholders
        if i >= len(subtitle_placeholders):
            print(f"  No placeholder available for subtitle {i+1}")
            continue
            
        # Use the image{i}subtitle naming format (i=1,2,3)
        subtitle_field = f'image{i+1}subtitle'
        
        # Get subtitle text
        subtitle_text = subtitles.get(subtitle_field, "")
        if not subtitle_text:
            print(f"  No text found for {subtitle_field}")
            continue
            
        try:
            shape = subtitle_placeholders[i]
            if hasattr(shape, "text_frame"):
                # Clear existing text
                shape.text_frame.clear()
                
                # Add new text
                p = shape.text_frame.paragraphs[0]
                p.text = subtitle_text
                print(f"  Added {subtitle_field}: {subtitle_text}")
            else:
                print(f"  Placeholder for {subtitle_field} has no text_frame")
        except Exception as e:
            print(f"  Error adding {subtitle_field}: {str(e)}")
    
    return slide_obj 

def create_content_with_logos_slide(prs, content_layout, slide_title, content, logo_paths):
    """
    Create a slide with content and up to three company/product logos.
    
    Args:
        prs: The presentation object
        content_layout: The slide layout to use
        slide_title: The slide title
        content: List of content items (bullet points)
        logo_paths: Dictionary with paths to logo images (logo1, logo2, logo3)
        
    Returns:
        The created content with logos slide
    """
    print(f"\n===== Creating ContentWithLogos slide: {slide_title} =====")
    print(f"LOGO DEBUG: Logo paths received: {logo_paths}")
    
    # Verify each logo file exists and print its details
    for logo_key, logo_path in logo_paths.items():
        if logo_path:
            file_exists = os.path.exists(logo_path)
            file_size = os.path.getsize(logo_path) if file_exists else 0
            print(f"LOGO DEBUG: {logo_key} -> Path: {logo_path}, Exists: {file_exists}, Size: {file_size} bytes")
        else:
            print(f"LOGO DEBUG: {logo_key} -> No path provided")
    
    slide_obj = prs.slides.add_slide(content_layout)
    
    # Print all shapes for debugging
    print("Available shapes in ContentWithLogos slide:")
    for i, shape in enumerate(slide_obj.shapes):
        shape_name = getattr(shape, "name", f"Shape {i}")
        shape_type = getattr(shape, "shape_type", "Unknown")
        placeholder_type = "Unknown"
        if hasattr(shape, "placeholder_format"):
            try:
                placeholder_type = shape.placeholder_format.type
            except:
                placeholder_type = "Error accessing type"
        shape_text = ""
        if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
            shape_text = shape.text_frame.text
        print(f"  Shape {i}: Name='{shape_name}', Type={shape_type}, Placeholder Type={placeholder_type}, Text='{shape_text}'")
    
    # Find title and content shapes using exact names from layout_analysis.json
    title_shape = find_shape_by_name(slide_obj, "Title")
    content_shape = find_shape_by_name(slide_obj, "Content")
    
    # If not found by name, fall back to placeholder types
    if not title_shape:
        for shape in slide_obj.shapes:
            if hasattr(shape, "placeholder_format"):
                try:
                    if shape.placeholder_format.type == 1:  # TITLE
                        title_shape = shape
                        print(f"  Found title shape by placeholder type")
                        break
                except:
                    pass
    
    if not content_shape:
        for shape in slide_obj.shapes:
            if hasattr(shape, "placeholder_format"):
                try:
                    if shape.placeholder_format.type == 2:  # BODY
                        content_shape = shape
                        print(f"  Found content shape by placeholder type")
                        break
                except:
                    pass
    
    # Set the title
    if title_shape and hasattr(title_shape, "text_frame"):
        title_shape.text_frame.clear()
        p = title_shape.text_frame.paragraphs[0]
        
        # Apply markdown formatting to title
        parse_markdown_to_runs(p, slide_title)
        print(f"  Set slide title to: {slide_title}")
        
        # Adjust title size based on length
        adjust_text_size(title_shape.text_frame, default_size=28, min_size=20, max_length=40, long_text_size=24, is_title=True)
    else:
        print("⚠️ WARNING: Could not find title shape for ContentWithLogos slide")
    
    # Set the content - use direct text setting instead of markdown formatting
    if content_shape and hasattr(content_shape, "text_frame"):
        # Clear the content placeholder
        content_shape.text_frame.clear()
        
        # Add each content item as a separate paragraph (without markdown formatting)
        for i, item in enumerate(content):
            if i == 0:
                # Use the first paragraph
                p = content_shape.text_frame.paragraphs[0]
            else:
                # Add a new paragraph for each subsequent item
                p = content_shape.text_frame.add_paragraph()
            
            # Set the paragraph text (no bold by default)
            p.text = item
            # Set bullet level to 0 (top level)
            p.level = 0
            # Ensure font is not bold by default
            if p.runs:
                p.runs[0].font.bold = False
        
        print(f"  Added {len(content)} content points to ContentWithLogos slide")
    else:
        print("⚠️ WARNING: Could not find content shape for ContentWithLogos slide")
    
    # Find logo placeholders by exact name from layout_analysis.json
    logo_placeholders = {
        "logo1": find_shape_by_name(slide_obj, "Logo1"),
        "logo2": find_shape_by_name(slide_obj, "Logo2"),
        "logo3": find_shape_by_name(slide_obj, "Logo3")
    }
    
    # Debug - list all found logo placeholders
    print(f"LOGO DEBUG: Found logo placeholders:")
    for key, placeholder in logo_placeholders.items():
        if placeholder:
            print(f"  - {key}: {getattr(placeholder, 'name', 'Unknown')}")
        else:
            print(f"  - {key}: Not found")
    
    # Try to find picture placeholders by type
    picture_placeholders = []
    for shape in slide_obj.shapes:
        try:
            if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 18:  # PICTURE
                picture_placeholders.append(shape)
                print(f"LOGO DEBUG: Found picture placeholder by type: {getattr(shape, 'name', 'Unknown')}")
        except Exception as e:
            print(f"LOGO DEBUG: Error checking placeholder type: {str(e)}")
    
    # Count of valid logo paths
    valid_logo_paths = {k: v for k, v in logo_paths.items() if v and os.path.exists(v)}
    valid_logo_count = len(valid_logo_paths)
    
    if valid_logo_count > 0:
        print(f"LOGO DEBUG: Adding {valid_logo_count} logos to slide")
        print(f"LOGO DEBUG: Valid logo paths:")
        for k, v in valid_logo_paths.items():
            file_size = os.path.getsize(v) if os.path.exists(v) else 0
            print(f"    - {k}: {v} (exists: {os.path.exists(v)}, size: {file_size} bytes)")
        
        # Match logos to placeholders
        for i, (logo_key, logo_path) in enumerate(valid_logo_paths.items()):
            logo_key_lower = logo_key.lower()
            placeholder = logo_placeholders.get(logo_key_lower)
            
            # Use picture placeholders if exact named placeholders aren't found
            if not placeholder and i < len(picture_placeholders):
                placeholder = picture_placeholders[i]
                print(f"LOGO DEBUG: Using generic picture placeholder {i} for {logo_key} (no named placeholder found)")
            
            if placeholder:
                print(f"LOGO DEBUG: Processing logo {logo_key} with placeholder {getattr(placeholder, 'name', 'Unknown')}")
                try:
                    # Get placeholder position and size
                    left = placeholder.left
                    top = placeholder.top
                    width = placeholder.width
                    height = placeholder.height
                    
                    print(f"LOGO DEBUG: Placeholder dimensions - Left: {left}, Top: {top}, Width: {width}, Height: {height}")
                    
                    # Try to validate the image file
                    try:
                        from PIL import Image
                        img = Image.open(logo_path)
                        print(f"LOGO DEBUG: Successfully opened image with PIL - Format: {img.format}, Size: {img.size}")
                        img.close()
                    except Exception as e:
                        print(f"LOGO DEBUG: Warning - Failed to validate image with PIL: {str(e)}")
                    
                    # First try inserting directly if the shape supports it
                    try:
                        if hasattr(placeholder, "insert_picture") and callable(placeholder.insert_picture):
                            placeholder.insert_picture(logo_path)
                            print(f"LOGO DEBUG: Added logo using insert_picture: {logo_path}")
                            continue  # Skip the rest of this logo's processing
                    except Exception as e:
                        print(f"LOGO DEBUG: insert_picture failed, trying alternative method: {str(e)}")
                    
                    # Convert SVG to PNG if it's an SVG file
                    if logo_path.lower().endswith('.svg'):
                        try:
                            import cairosvg
                            import tempfile
                            
                            # Create a temporary PNG file
                            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                                tmp_path = tmp.name
                            
                            # Convert SVG to PNG
                            cairosvg.svg2png(url=logo_path, write_to=tmp_path)
                            print(f"LOGO DEBUG: Converted SVG to PNG: {tmp_path}")
                            
                            # Use the PNG file instead
                            logo_path = tmp_path
                        except ImportError:
                            print("LOGO DEBUG: cairosvg not available, using SVG directly")
                        except Exception as e:
                            print(f"LOGO DEBUG: Error converting SVG to PNG: {str(e)}")
                    
                    # Delete the placeholder
                    try:
                        placeholder._element.getparent().remove(placeholder._element)
                        print(f"LOGO DEBUG: Successfully removed placeholder element")
                    except Exception as e:
                        print(f"LOGO DEBUG: Error removing placeholder (continuing anyway): {str(e)}")
                    
                    # Add the logo picture at the same position
                    try:
                        picture = slide_obj.shapes.add_picture(
                            logo_path,
                            left,
                            top,
                            width=width,
                            height=height
                        )
                        print(f"LOGO DEBUG: Successfully added {logo_key} logo - Shape ID: {id(picture)}")
                    except Exception as e:
                        print(f"LOGO DEBUG: ❌ ERROR: Failed to add picture: {str(e)}")
                        # Try to get more details about the error
                        import traceback
                        print(f"LOGO DEBUG: Error details:\n{traceback.format_exc()}")
                        
                except Exception as e:
                    print(f"LOGO DEBUG: ❌ ERROR: Failed to process {logo_key} logo: {str(e)}")
                    # Try fallback method - add without using placeholder
                    try:
                        # Calculate a position at the bottom of the slide
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height
                        
                        # Logo dimensions
                        logo_width = Inches(1.5)
                        logo_height = Inches(1.5)
                        
                        # Calculate position (centered horizontally, near bottom of slide)
                        x_position = (slide_width - logo_width) / 2
                        y_position = slide_height - logo_height - Inches(0.5)
                        
                        picture = slide_obj.shapes.add_picture(
                            logo_path,
                            x_position,
                            y_position,
                            width=logo_width,
                            height=logo_height
                        )
                        print(f"LOGO DEBUG: Added {logo_key} using fallback positioning method")
                    except Exception as e2:
                        print(f"LOGO DEBUG: ❌ Fallback also failed: {str(e2)}")
            else:
                # No matching placeholder found, add at a calculated position
                try:
                    print(f"LOGO DEBUG: No placeholder found for {logo_key}, using calculated position")
                    # Calculate default position
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    logo_width = Inches(1.5)
                    logo_height = Inches(1.5)
                    
                    # Position logos at bottom, spaced evenly
                    logo_num = int(logo_key.replace("logo", ""))
                    total_logos = len(valid_logo_paths)
                    
                    # Calculate logo position and spacing
                    total_width = total_logos * logo_width
                    spacing = (slide_width - total_width) / (total_logos + 1)
                    x_position = spacing * logo_num + (logo_num - 1) * logo_width
                    y_position = slide_height - logo_height - Inches(0.75)
                    
                    picture = slide_obj.shapes.add_picture(
                        logo_path,
                        x_position,
                        y_position,
                        width=logo_width,
                        height=logo_height
                    )
                    print(f"LOGO DEBUG: Successfully added {logo_key} at calculated position")
                except Exception as e:
                    print(f"LOGO DEBUG: ❌ ERROR: Failed to add {logo_key} at calculated position: {str(e)}")
    else:
        print("LOGO DEBUG: No valid logo paths provided, creating slide without logos")
    
    # Final check - count shapes after adding logos
    print(f"LOGO DEBUG: Final shape count: {len(list(slide_obj.shapes))}")
    
    return slide_obj 