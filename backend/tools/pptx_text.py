"""
Functions for handling text formatting in PowerPoint slides
"""
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

def adjust_text_size(text_frame, default_size=24, min_size=14, max_length=50, long_text_size=16, is_title=False):
    """
    Adjust the text size based on the length of the text content.
    Makes text smaller for longer content and ensures proper wrapping.
    """
    if not text_frame or not hasattr(text_frame, "text") or not text_frame.text:
        return
    
    text_length = len(text_frame.text)
    text_preview = text_frame.text[:50] + ("..." if len(text_frame.text) > 50 else "")
    print(f"Adjusting text size for: '{text_preview}' length={text_length}")
    
    # Clear any auto-size settings to allow manual size control
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    
    # Apply word wrapping to ensure text fits within boundaries
    text_frame.word_wrap = True
    
    # Calculate appropriate font size based on text length
    font_size = default_size
    if text_length > max_length * 1.5:  # Very long text
        font_size = min_size
        print(f"  Setting to minimum size for very long text: {min_size}pt")
    elif text_length > max_length:  # Long text
        font_size = long_text_size
        print(f"  Setting to smaller size for long text: {long_text_size}pt")
    else:
        print(f"  Setting to default size: {default_size}pt")
    
    # Set the font size for all paragraphs
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        
        # For titles, make bold but do not center
        if is_title and paragraph.runs:
            paragraph.runs[0].font.bold = True
        
        # Make sure text wraps properly
        paragraph.word_wrap = True
        
    # If there are multiple lines, reduce the spacing between paragraphs
    if len(text_frame.paragraphs) > 1:
        for i, paragraph in enumerate(text_frame.paragraphs):
            # Skip the first paragraph as it doesn't need spacing before
            if i > 0:
                paragraph.space_before = Pt(2)  # Minimal space between paragraphs

def format_content_text(content_shape, content_items):
    """Format content text with appropriate styles and sizing"""
    if not content_shape or not hasattr(content_shape, "text_frame"):
        return
    
    # Clear the content placeholder
    content_shape.text_frame.clear()
    
    # Apply word wrapping
    content_shape.text_frame.word_wrap = True
    
    # Determine appropriate font size based on total content length and item count
    total_length = sum(len(item) for item in content_items)
    num_items = len(content_items)
    
    print(f"Formatting content with {num_items} items, total length: {total_length}")
    
    # Determine base font size based on content volume
    base_size = 18  # Default size
    if total_length > 800 or num_items > 10:
        base_size = 12  # Very dense content
        print(f"  Using small font (12pt) for very dense content")
    elif total_length > 500 or num_items > 7:
        base_size = 14  # Dense content
        print(f"  Using reduced font (14pt) for dense content")
    elif total_length > 300 or num_items > 5:
        base_size = 16  # Moderate content
        print(f"  Using medium font (16pt) for moderate content")
    else:
        print(f"  Using standard font (18pt) for normal content")
    
    # Add each content item as a paragraph
    for idx, item in enumerate(content_items):
        if idx == 0:
            p = content_shape.text_frame.paragraphs[0]
        else:
            p = content_shape.text_frame.add_paragraph()
        
        p.text = item
        p.level = 0
        p.font.size = Pt(base_size)
        
        # Add a bit of spacing between items but not too much
        if idx > 0:
            p.space_before = Pt(6)
    
    print(f"  Added {len(content_items)} formatted content items") 