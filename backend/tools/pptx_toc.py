"""
Functions for handling PowerPoint Table of Contents slides
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .pptx_shapes import find_shape_by_name
from .pptx_text import adjust_text_size
from .pptx_markdown import parse_markdown_to_runs
  
def _find_section_text_shapes(slide):
    """
    Find shapes that correspond to sections in the Table of Contents.
    In the master slide layout, these are named "Section1", "Section2", etc.
    In actual slides, they are typically "Text Placeholder X" in order (2-9).
    
    Args:
        slide: The slide object to search for section shapes
        
    Returns:
        dict: Dictionary mapping section numbers (int) to shape objects
    """
    section_shapes = {}
    
    # First try by name (for master layouts)
    for i in range(1, 9):  # Sections 1-8
        section_name = f"Section{i}"
        section_shape = find_shape_by_name(slide, section_name)
        
        if section_shape:
            section_shapes[i] = section_shape
    
    # If we found section shapes by name, return them
    if section_shapes:
        return section_shapes
    
    # Otherwise, look for text placeholders for sections
    # In a created slide, they are typically Text Placeholder 2-9
    text_placeholders = []
    for shape in slide.shapes:
        if hasattr(shape, "name") and "Text Placeholder" in shape.name and hasattr(shape, "text_frame"):
            text_placeholders.append(shape)
    
    # Sort placeholders by index number in name
    try:
        text_placeholders.sort(key=lambda s: int(s.name.replace("Text Placeholder ", "")))
    except:
        # If sorting fails, leave them in the order they were found
        pass
    
    # Map to section numbers
    for i, shape in enumerate(text_placeholders):
        # Use 1-based indexing for section numbers
        section_num = i + 1
        if section_num <= 8:  # Only sections 1-8
            section_shapes[section_num] = shape
    
    return section_shapes

def create_table_of_contents_slide(prs, toc_data):
    """
    Create a Table of Contents slide using the template layout.
    
    Args:
        prs: The presentation object
        toc_data: Dictionary with 'title' and 'sections' keys
        
    Returns:
        The created TOC slide
    """
    # Use the TableOfContents layout
    for idx, layout in enumerate(prs.slide_layouts):
        layout_name = getattr(layout, 'name', f'Layout {idx}')
        if 'tableofcontents' in layout_name.lower().replace(' ', ''):
            toc_layout = layout
            break
    else:
        # Fallback to index 12 which is typically TableOfContents
        toc_layout = prs.slide_layouts[12] if len(prs.slide_layouts) > 12 else prs.slide_layouts[0]
    
    # Create the slide with the TOC layout
    toc_slide = prs.slides.add_slide(toc_layout)
    
    # # Set the title
    # title_shape = _find_title_shape(toc_slide)
    # if title_shape and hasattr(title_shape, 'text_frame'):
    #     if len(title_shape.text_frame.paragraphs) > 0:
    #         _preserve_text_style(title_shape.text_frame.paragraphs[0], "Table of Contents")
    #     else:
    #         title_shape.text_frame.text = "Table of Contents"
    
    # Get the section titles
    section_titles = toc_data.get('sections', [])
    section_count = min(len(section_titles), 8)  # Max 8 sections
    
    # Find section shapes
    section_shapes = _find_section_text_shapes(toc_slide)
    
    # Set section text in placeholders
    for i in range(1, section_count + 1):
        if i in section_shapes and i <= len(section_titles):
            section_shape = section_shapes[i]
            section_title = section_titles[i-1]  # Adjust for 0-based index in section_titles
            
            if hasattr(section_shape, 'text_frame'):
                # Clear the shape text frame
                section_shape.text_frame.clear()
                
                # Get the first paragraph
                p = section_shape.text_frame.paragraphs[0]
                
                # Apply markdown formatting to section title
                parse_markdown_to_runs(p, section_title)
                adjust_text_size(section_shape.text_frame, default_size=24, min_size=18, max_length=60, long_text_size=20)
    
    return toc_slide

def process_toc_slide(toc_slide, section_titles, max_sections=8):
    """
    Process a Table of Contents slide by updating section titles and formatting.
    
    Args:
        toc_slide: The TOC slide object
        section_titles: List of section titles to display
        max_sections: Maximum number of sections supported in the template
    """
    
    # Find section shapes
    section_shapes = _find_section_text_shapes(toc_slide)
    
    # Process each section
    section_count = min(len(section_titles), max_sections)
    
    for i in range(1, section_count + 1):
        if i in section_shapes and i <= len(section_titles):
            section_shape = section_shapes[i]
            section_title = section_titles[i-1]  # Adjust for 0-based index in section_titles
            
            if hasattr(section_shape, 'text_frame'):
                # Clear the shape text frame
                section_shape.text_frame.clear()
                
                # Get the first paragraph
                p = section_shape.text_frame.paragraphs[0]
                
                # Apply markdown formatting to section title
                parse_markdown_to_runs(p, section_title)
    
    # Hide unused section shapes
    for i in range(section_count + 1, 9):
        if i in section_shapes:
            section_shape = section_shapes[i]
            try:
                section_shape.element.getparent().remove(section_shape.element)
            except:
                try:
                    section_shape.left = -10000
                    section_shape.top = -10000
                except:
                    pass
    
    return toc_slide 