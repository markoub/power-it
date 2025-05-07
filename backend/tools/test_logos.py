import os
import sys
import asyncio
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from tools.generate_pptx import generate_pptx_from_slides
from pptx import Presentation

async def test_content_with_logos():
    """Test function to diagnose ContentWithLogos issues"""
    print("\n=== Starting ContentWithLogos diagnostic test ===\n")
    
    # Create test slide data
    slides = [
        {
            "id": 0,
            "title": "Test Presentation",
            "type": "Welcome",
            "fields": {
                "subtitle": "Demo presentation for logo testing"
            }
        },
        {
            "id": 1,
            "title": "Test Content With Logos",
            "type": "ContentWithLogos",
            "fields": {
                "content": [
                    "This is a test slide with logos",
                    "The logos should appear below",
                    "If they don't appear, there's an issue"
                ],
                "logo1": "Amazon AWS",
                "logo2": "Google Cloud",
                "logo3": "Microsoft Azure"
            }
        }
    ]
    
    print("Test slides data:")
    for slide in slides:
        print(f"  Slide ID {slide.get('id')}: type={slide.get('type')}, title={slide.get('title')}")
        if 'fields' in slide and slide['fields']:
            print(f"    Fields: {slide['fields'].keys()}")
            if 'content' in slide['fields']:
                print(f"    Content: {slide['fields']['content']}")
    
    # Create test presentation in current directory
    output_path = "logo_test.pptx"
    
    # Generate the presentation
    try:
        # Modify the save path in generate_pptx to save directly
        old_pptx_path = None
        
        result = await generate_pptx_from_slides(slides, output_path)
        print(f"Generation result: {result}")
        
        # Copy from storage to local directory
        if os.path.exists(result.pptx_path):
            import shutil
            shutil.copy(result.pptx_path, output_path)
            print(f"Copied presentation from {result.pptx_path} to {output_path}")
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        import traceback
        print(traceback.format_exc())
    
    # Verify the output
    if os.path.exists(output_path):
        print(f"\n✅ Test presentation created successfully: {output_path}")
        # Open the presentation to verify logos
        prs = Presentation(output_path)
        
        # Check slide count
        print(f"Presentation has {len(prs.slides)} slides")
        
        # Check each slide for ContentWithLogos
        content_with_logos_slide = None
        for i, slide in enumerate(prs.slides):
            # Try to find content with slide layout name
            layout_name = "Unknown"
            if hasattr(slide, "slide_layout") and hasattr(slide.slide_layout, "name"):
                layout_name = slide.slide_layout.name
            
            print(f"Slide {i+1} uses layout: {layout_name}")
            
            if layout_name == "ContentWithLogos":
                content_with_logos_slide = slide
                print(f"Found ContentWithLogos slide at position {i+1}")
                break
        
        # Analyze ContentWithLogos slide if found
        if content_with_logos_slide:
            # Count the shapes
            print(f"ContentWithLogos slide has {len(content_with_logos_slide.shapes)} shapes")
            
            # List all shapes
            print("\nShapes on ContentWithLogos slide:")
            for i, shape in enumerate(content_with_logos_slide.shapes):
                shape_type = getattr(shape, "shape_type", "Unknown")
                if hasattr(shape, "shape_type"):
                    shape_type_name = str(shape.shape_type).split('.')[-1]
                else:
                    shape_type_name = "Unknown"
                
                # Get shape name
                shape_name = getattr(shape, "name", f"Shape {i}")
                
                # Get text if available
                shape_text = ""
                if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                    shape_text = shape.text_frame.text[:50] + "..." if len(shape.text_frame.text) > 50 else shape.text_frame.text
                
                print(f"  {i}: {shape_name} (Type: {shape_type_name}) - Text: '{shape_text}'")
        else:
            print("❌ No ContentWithLogos slide found in the presentation!")
    else:
        print(f"\n❌ Failed to create test presentation: {output_path}")
    
    print("\n=== ContentWithLogos diagnostic test completed ===\n")
    
if __name__ == "__main__":
    asyncio.run(test_content_with_logos()) 