from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_test_toc_slide():
    """Test creating a table of contents slide with the fixed method"""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Use title and content layout
    
    # Create a slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = "Table of Contents"
    
    # Get content placeholder
    content_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # Content
            content_shape = shape
            break
    
    # If no content shape found, add a text box
    if not content_shape:
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        content_shape = slide.shapes.add_textbox(left, top, width, height)
    
    # Sample section titles
    section_titles = [
        "Understanding Eyelid Aging",
        "Surgical Approaches",
        "Non-Surgical Options",
        "Treatment Algorithms"
    ]
    
    # Create formatted TOC with numbered items
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for i, section_title in enumerate(section_titles, 1):
        if i == 1:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        
        # Use a numbered format instead of bullets
        paragraph.text = f"{i:02d}  {section_title}"
        paragraph.level = 0
        paragraph.font.size = Pt(24)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.LEFT
        
        # Add spacing between items
        if i > 1:
            paragraph.space_before = Pt(12)
    
    # Save the presentation
    prs.save("test_toc.pptx")
    print("Created test_toc.pptx with numbered TOC items")

if __name__ == "__main__":
    create_test_toc_slide() 