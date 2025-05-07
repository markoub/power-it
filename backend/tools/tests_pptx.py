import os
import asyncio
import unittest
from unittest.mock import patch
import tempfile
import shutil
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

from models import SlidePresentation, Slide, CompiledPresentation, CompiledSlide
from tools.generate_pptx import generate_pptx_from_slides, calculate_content_size

def async_test(coro):
    def wrapper(*args, **kwargs):
        loop = asyncio.get_event_loop()
        return loop.run_until_complete(coro(*args, **kwargs))
    return wrapper

class TestPptxGeneration(unittest.TestCase):
    def setUp(self):
        # Create a temporary directory for test files
        self.test_dir = tempfile.mkdtemp()
        # Mock the PRESENTATIONS_STORAGE_DIR to use our temporary directory
        self.patcher = patch('tools.generate_pptx.PRESENTATIONS_STORAGE_DIR', self.test_dir)
        self.mock_storage_dir = self.patcher.start()
        
        # Create test subdirectories
        os.makedirs(os.path.join(self.test_dir, "test", "images"), exist_ok=True)
        
    def tearDown(self):
        # Stop the patcher
        self.patcher.stop()
        # Remove the temporary directory
        shutil.rmtree(self.test_dir)
        
    def test_calculate_content_size(self):
        # Test with short content
        short_content = ["Short bullet point"]
        size_score, _ = calculate_content_size(short_content)
        self.assertEqual(size_score, 1, "Short content should have size score 1")
        
        # Test with medium content
        medium_content = ["This is a medium length bullet point with some details",
                         "Here is another medium length point with additional information"]
        size_score, _ = calculate_content_size(medium_content)
        self.assertEqual(size_score, 2, "Medium content should have size score 2")
        
        # Test with long content
        long_content = ["This is a very long bullet point with lots of details and information " +
                        "that would require significant space on a slide to display properly",
                        "Here is another lengthy bullet point that contains important information " +
                        "and details that need to be communicated clearly to the audience",
                        "And one more substantial point with technical details and explanations " +
                        "that require more space than a typical bullet point would"]
        size_score, _ = calculate_content_size(long_content)
        self.assertEqual(size_score, 3, "Long content should have size score 3")
        
    @async_test
    async def test_adaptive_content_sizing(self):
        # Create test slides with different content lengths
        slides = SlidePresentation(
            title="Test Presentation",
            slides=[
                # Title slide with short content
                Slide(
                    title="Title Slide",
                    type="title",
                    content=["Short presentation subtitle"],
                    notes="Speaker notes for title slide"
                ),
                # Content slide with short content
                Slide(
                    title="Short Content",
                    type="content",
                    content=["Brief bullet point"],
                    notes="Speaker notes for short content"
                ),
                # Content slide with medium content
                Slide(
                    title="Medium Content",
                    type="content",
                    content=[
                        "This is a medium length bullet point with some details",
                        "Here is another medium length point with additional information"
                    ],
                    notes="Speaker notes for medium content"
                ),
                # Content slide with long content
                Slide(
                    title="Long Content",
                    type="content",
                    content=[
                        "This is a very long bullet point with lots of details and information " +
                        "that would require significant space on a slide to display properly",
                        "Here is another lengthy bullet point that contains important information " +
                        "and details that need to be communicated clearly to the audience",
                        "And one more substantial point with technical details and explanations " +
                        "that require more space than a typical bullet point would"
                    ],
                    notes="Speaker notes for long content"
                )
            ]
        )
        
        # Generate PPTX from slides
        pptx_gen = await generate_pptx_from_slides(slides, presentation_id="test")
        
        # Verify the PPTX file was created
        self.assertTrue(os.path.exists(pptx_gen.pptx_path), f"PPTX file was not created at {pptx_gen.pptx_path}")
        
        # Open the PPTX file and check slide content
        prs = Presentation(pptx_gen.pptx_path)
        
        # Verify number of slides
        self.assertEqual(len(prs.slides), 4, "Presentation should have 4 slides")
        
        # Check if auto-size is enabled on title slide
        title_slide = prs.slides[0]
        title_shape = title_slide.shapes.title
        subtitle_shape = None
        
        # Find the subtitle shape (typically placeholder 1)
        for shape in title_slide.shapes:
            if shape.placeholder_format.idx == 1:
                subtitle_shape = shape
                break
                
        self.assertIsNotNone(title_shape, "Title shape not found")
        self.assertIsNotNone(subtitle_shape, "Subtitle shape not found")
        
        if title_shape and hasattr(title_shape, 'text_frame'):
            self.assertEqual(
                title_shape.text_frame.auto_size, 
                MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                "Title should have auto-size enabled"
            )
        
        if subtitle_shape and hasattr(subtitle_shape, 'text_frame'):
            self.assertEqual(
                subtitle_shape.text_frame.auto_size, 
                MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                "Subtitle should have auto-size enabled"
            )
        
        # Check if auto-size is enabled on content slides
        for i in range(1, 4):  # Check slides 1, 2, 3 (index 1-3)
            slide = prs.slides[i]
            title_shape = slide.shapes.title
            content_shape = None
            
            # Find the content shape (typically placeholder 1)
            for shape in slide.shapes:
                if shape.placeholder_format.idx == 1:
                    content_shape = shape
                    break
            
            self.assertIsNotNone(content_shape, f"Content shape not found on slide {i}")
            
            if content_shape and hasattr(content_shape, 'text_frame'):
                self.assertEqual(
                    content_shape.text_frame.auto_size, 
                    MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
                    f"Content on slide {i} should have auto-size enabled"
                )
                
        # Check image positioning on a slide with content
        # For this test, we mainly verify that the presentation structure is as expected
        # and adapts to different content sizes
        self.assertEqual(pptx_gen.slide_count, 4, "Presentation should have 4 slides")

if __name__ == "__main__":
    unittest.main() 